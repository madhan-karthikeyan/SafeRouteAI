#!/usr/bin/env python3
import sys
from pptx import Presentation

def inspect_template():
    pptx_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/IDEA_Presentation_Format.pptx"
    prs = Presentation(pptx_path)
    
    print(f"=== TEMPLATE INSPECTION: {pptx_path} ===")
    print(f"Slide Width: {prs.slide_width.inches} inches")
    print(f"Slide Height: {prs.slide_height.inches} inches")
    print(f"Total Slides in Template: {len(prs.slides)}\n")
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"--- SLIDE {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                print(f"  [Text]: {shape.text_frame.text.strip().replace('\n', ' | ')}")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            print(f"  [Notes]: {slide.notes_slide.notes_text_frame.text.strip().replace('\n', ' | ')}")
        print()

if __name__ == "__main__":
    inspect_template()
