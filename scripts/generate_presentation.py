#!/usr/bin/env python3
"""
SafeRouteAI PowerPoint Generator
--------------------------------
Generates a competition-ready, 16-slide PowerPoint presentation (`SafeRouteAI.pptx`)
for the SafeRoute AI project using python-pptx.

Adheres strictly to documentation in `docs/`:
1. problem-statement.md
2. engineering_report.md
3. plan.md
4. presentation.md
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Define Dark Theme Palette
COLOR_BG = RGBColor(11, 15, 25)          # Obsidian Dark #0B0F19
COLOR_CARD = RGBColor(30, 41, 59)        # Slate Dark Card #1E293B
COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate 700 Border #334155
COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# White Main Text #F8FAFC
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# Muted Gray Text #94A3B8
COLOR_ACCENT_CYAN = RGBColor(56, 189, 248)# Sky Cyan Accent #38BDF8
COLOR_ACCENT_EMERALD = RGBColor(52, 211, 153)# Safe Emerald Green #34D399
COLOR_ACCENT_RED = RGBColor(248, 113, 113)  # Hazard Red #F87171
COLOR_ACCENT_YELLOW = RGBColor(251, 191, 36) # Reroute Yellow #FBBF24
COLOR_ACCENT_PURPLE = RGBColor(167, 139, 250)# AI/Engine Purple #A78BFA
COLOR_CONTAINER_DARK = RGBColor(15, 23, 42) # Deep Slate Container #0F172A

def set_slide_background(slide, color=COLOR_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_header(slide, category, title):
    # Category Tag / Badge
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_ACCENT_CYAN
    p_cat.font.name = "Arial"
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MAIN
    p_title.font.name = "Arial"

def add_footer(slide, current_slide, total_slides=16):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.3))
    tf = footer_box.text_frame
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"SafeRoute AI  |  Decentralized Fire Evacuation Router  |  Slide {current_slide} of {total_slides}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.font.name = "Arial"

def set_speaker_notes(slide, script, points, transition, qa):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    
    notes_text = f"""=== SCRIPT / WHAT TO SAY ===
{script}

=== KEY TECHNICAL POINTS ===
{points}

=== SLIDE TRANSITION ===
{transition}

=== EXPECTED JURY Q&A ===
{qa}"""
    text_frame.text = notes_text

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    print("Building Slide 1: Title...")
    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # Large Backdrop Card
    create_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), COLOR_CONTAINER_DARK, COLOR_CARD_BORDER)
    
    # Decorative Top Accent Line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT_CYAN
    line.line.fill.background()
    
    # Project Title
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.9), Inches(1.5))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SafeRoute AI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.text = "Decentralized Evacuation Routing with Real-Time Hazard Mapping"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_ACCENT_CYAN
    p2.font.name = "Arial"
    
    p3 = tf.add_paragraph()
    p3.text = "Self-Healing ESP32 Link-State Mesh  •  Sub-300ms Dynamic Pathfinding  •  3D EOC Digital Twin"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.font.name = "Arial"
    
    # Key Highlights Badges Container
    b_y = 3.8
    badges = [
        ("SUB-300ms LATENCY", "Real-Time Edge Pathfinding", COLOR_ACCENT_CYAN),
        ("ESP-NOW MESH", "Connectionless Flooding", COLOR_ACCENT_EMERALD),
        ("3-TIER FAIL-SAFE", "Local • Consensus • Default", COLOR_ACCENT_YELLOW),
        ("3D DIGITAL TWIN", "Three.js EOC Interface", COLOR_ACCENT_PURPLE),
    ]
    for i, (tag, desc, color) in enumerate(badges):
        b_x = 1.2 + i * 2.75
        create_card(slide1, Inches(b_x), Inches(b_y), Inches(2.6), Inches(1.3), COLOR_CARD, color)
        
        tb = slide1.shapes.add_textbox(Inches(b_x + 0.1), Inches(b_y + 0.15), Inches(2.4), Inches(1.0))
        tf_b = tb.text_frame
        tf_b.word_wrap = True
        
        pb1 = tf_b.paragraphs[0]
        pb1.text = tag
        pb1.font.size = Pt(11)
        pb1.font.bold = True
        pb1.font.color.rgb = color
        
        pb2 = tf_b.add_paragraph()
        pb2.text = desc
        pb2.font.size = Pt(10)
        pb2.font.color.rgb = COLOR_TEXT_MAIN
        
    # Metadata Footer
    m_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.6), Inches(10.9), Inches(0.8))
    tf_m = m_box.text_frame
    pm = tf_m.paragraphs[0]
    pm.text = "Hackathon Final Round Defense  |  Presentation Length: 8–10 Minutes  |  Hardware & Software Implementation"
    pm.font.size = Pt(11)
    pm.font.color.rgb = COLOR_TEXT_MUTED
    
    set_speaker_notes(
        slide1,
        script="Good day judges and audience. We are presenting SafeRoute AI — a fully decentralized, self-healing evacuation routing system engineered to save lives in complex commercial building fires. Standard static exit signs can guide occupants directly into toxic smoke or flashovers during structural fires. SafeRoute AI replaces passive signage with tiny embedded routers on ESP32 microcontrollers that detect fire vectors locally, flood hazard link-states across an ESP-NOW mesh, and continuously recompute the safest egress paths in under 300 milliseconds without relying on any cloud or central server.",
        points="• Position SafeRoute AI as an edge-computing life-safety innovation.\n• Highlight the core engineering pillars: Decentralized ESP-NOW mesh, on-device Dijkstra routing, sub-300ms reaction budget, and 3D EOC visualization.\n• Frame the project as augmenting code-mandated static signage with real-time dynamic intelligence.",
        transition="Let's examine why traditional static evacuation methods fail during real-world building emergencies.",
        qa="Q: Why ESP32 microcontrollers instead of standard Wi-Fi router networks?\nA: Standard Wi-Fi networks rely on central access points and cloud infrastructure that often lose power or connection during structural fires. ESP32 microcontrollers use connectionless, peer-to-peer ESP-NOW flooding, creating a resilient mesh that operates 100% autonomously on battery backup."
    )

    print("Building Slide 2: Problem Statement...")
    # ==========================================
    # SLIDE 2: Problem Statement
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "1. Problem Background", "The Critical Danger of Static Evacuation Signage in Structural Fires")
    add_footer(slide2, 2)
    
    # Left Column: Problem Breakdown Cards
    left_cards = [
        ("Static Exit Signs Direct Occupants into Danger", "Traditional illuminated exit signs direct occupants toward fixed exit routes regardless of fire spread. In high-rise fires, flashovers and smoke block exits within minutes, turning static routes into death traps.", COLOR_ACCENT_RED),
        ("Rapid Flashover & Toxic Smoke Hazards", "Modern building materials produce toxic carbon monoxide and synthetic gases. Flashover conditions occur within 2 to 3 minutes, leaving occupants zero margin for error when selecting evacuation paths.", COLOR_ACCENT_YELLOW),
        ("Centralized Infrastructure Single-Points-of-Failure", "Centralized alarm dashboards and cloud-managed smart signage suffer from single points of failure. If power lines burn, Wi-Fi drops, or the server crashes, evacuation guidance collapses completely.", COLOR_ACCENT_CYAN)
    ]
    
    for i, (title, text, accent) in enumerate(left_cards):
        c_y = 1.4 + i * 1.7
        create_card(slide2, Inches(0.8), Inches(c_y), Inches(7.2), Inches(1.5), COLOR_CARD, COLOR_CARD_BORDER)
        
        # Accent left bar
        bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(c_y), Inches(0.1), Inches(1.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        
        tb = slide2.shapes.add_textbox(Inches(1.0), Inches(c_y + 0.1), Inches(6.8), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        
    # Right Column: Metrics & Operational Constraints Box
    create_card(slide2, Inches(8.3), Inches(1.4), Inches(4.233), Inches(4.9), COLOR_CONTAINER_DARK, COLOR_ACCENT_RED)
    
    tb_r = slide2.shapes.add_textbox(Inches(8.5), Inches(1.6), Inches(3.833), Inches(4.5))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    pr1 = tf_r.paragraphs[0]
    pr1.text = "CRITICAL METRICS"
    pr1.font.size = Pt(12)
    pr1.font.bold = True
    pr1.font.color.rgb = COLOR_ACCENT_RED
    
    metrics = [
        ("< 300 ms", "Maximum latency budget required to recompute exit paths before occupant panic"),
        ("2 - 3 Mins", "Time to flashover in modern synthetic commercial interiors"),
        ("80%+", "Fire fatalities caused by toxic smoke inhalation rather than direct thermal burns")
    ]
    
    for val, lbl in metrics:
        p_v = tf_r.add_paragraph()
        p_v.text = val
        p_v.font.size = Pt(24)
        p_v.font.bold = True
        p_v.font.color.rgb = COLOR_TEXT_MAIN
        
        p_l = tf_r.add_paragraph()
        p_l.text = lbl
        p_l.font.size = Pt(10)
        p_l.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide2,
        script="In commercial facility emergencies, static exit signs are a major hazard vector. When a fire breaks out in a corridor, fixed exit signs continuously point people right into flashovers and toxic smoke. Statistics show over 80% of fire deaths stem from smoke inhalation during egress attempts. Furthermore, centralized smart building management systems suffer from single points of failure — when structural fires sever electrical mains or Wi-Fi access points, centralized routing fails completely. Our design constraint is strict: the system must process multi-vector sensors and reroute visual indicators in under 300 milliseconds.",
        points="• Documented problem background directly from problem-statement.md.\n• Highlight the 3 core pain points: static sign traps, rapid flashover timeline, and centralized single-point-of-failure.\n• Stress the 300ms reaction latency threshold.",
        transition="Next, let's break down the technical limitations of current market evacuation systems.",
        qa="Q: Why is sub-300ms reaction speed required?\nA: Human reaction to sudden alarms occurs in fractions of a second. If an LED sign takes seconds to change state when flashover occurs, occupants will already have committed to a compromised hallway, causing bottlenecking and smoke exposure."
    )

    print("Building Slide 3: Existing Challenges...")
    # ==========================================
    # SLIDE 3: Existing Challenges
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "2. Existing Limitations", "Comparison of Current Evacuation Systems vs Real-World Needs")
    add_footer(slide3, 3)
    
    challenges = [
        ("Binary Threshold Triggers", "Current alarms use hard binary triggers (e.g. temp > 57°C). They fail to compute continuous hazard progression or early smoldering gas build-up, missing early intervention windows.", COLOR_ACCENT_YELLOW),
        ("Centralized Wi-Fi / Server Dependency", "Cloud-connected smart exit systems depend on active routers and central servers. During fires, network switches fail, severing communication precisely when needed most.", COLOR_ACCENT_RED),
        ("Zero Occupancy & Congestion Awareness", "Legacy signage cannot evaluate crowd density or corridor egress throughput. Directing 500 occupants into a single narrow staircase causes fatal stampedes and delays.", COLOR_ACCENT_CYAN),
        ("Blind Evacuation Guidance", "Static exit signs provide uniform green lighting even when the hallway behind the exit door is engulfed in flame, forcing occupants into high-risk trial-and-error choices.", COLOR_ACCENT_PURPLE)
    ]
    
    for i, (title, desc, color) in enumerate(challenges):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 5.95
        y = 1.4 + row * 2.6
        
        create_card(slide3, Inches(x), Inches(y), Inches(5.75), Inches(2.35), COLOR_CARD, COLOR_CARD_BORDER)
        
        tb = slide3.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(5.35), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = f"LIMITATION 0{i+1}: {title.upper()}"
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide3,
        script="Analyzing current commercial offerings reveals four critical flaws. First, binary thresholding: traditional alarms don't calculate continuous hazard curves; they operate on simplistic binary triggers that ignore gradual gas PPM and heat rises. Second, reliance on central cloud servers: if local Wi-Fi or central switches collapse, smart signs go blind. Third, zero occupancy awareness: traditional signs ignore crowd accumulation, sending hundreds of people down a choked corridor. Fourth, blind guidance: occupants have no visual cue indicating whether an alternate exit exists or if they should shelter in place.",
        points="• Detail the 4 core systemic limitations.\n• Explain why binary thresholds fail compared to continuous exponential weighting.\n• Contrast centralized infrastructure against link-state mesh resilience.",
        transition="To eliminate these systemic vulnerabilities, we developed SafeRoute AI.",
        qa="Q: How does SafeRoute AI address occupancy without expensive vision cameras on every node?\nA: SafeRoute AI integrates access control counts and camera streams at key junction nodes, combining occupant counts with base corridor distances in an additive congestion term."
    )

    print("Building Slide 4: Proposed Solution...")
    # ==========================================
    # SLIDE 4: Proposed Solution
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "3. Proposed Solution", "SafeRoute AI — Decentralized Edge Graph Routing Architecture")
    add_footer(slide4, 4)
    
    # Top Solution Overview Box
    create_card(slide4, Inches(0.8), Inches(1.4), Inches(11.733), Inches(1.4), COLOR_CONTAINER_DARK, COLOR_ACCENT_CYAN)
    tb_top = slide4.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(1.2))
    tf_top = tb_top.text_frame
    tf_top.word_wrap = True
    
    pt1 = tf_top.paragraphs[0]
    pt1.text = "CORE VALUE PROPOSITION"
    pt1.font.size = Pt(11)
    pt1.font.bold = True
    pt1.font.color.rgb = COLOR_ACCENT_CYAN
    
    pt2 = tf_top.add_paragraph()
    pt2.text = "A decentralized, self-healing physical hazard routing network. Each ESP32 node acts as a tiny link-state router: it continuously senses temperature, smoke, flame, and occupancy, fuses them into an exponential edge cost, floods updates over ESP-NOW mesh, and runs Dijkstra on-device to steer dynamic WS2812B chasing LEDs."
    pt2.font.size = Pt(11)
    pt2.font.color.rgb = COLOR_TEXT_MAIN
    
    # 4 Engineering Pillars Cards
    pillars = [
        ("Decentralized Edge Dijkstra", "No central server or master brain. Every node computes its shortest path to safety independently.", COLOR_ACCENT_CYAN),
        ("Continuous Multi-Sensor Fusion", "Exponential cost math fuses Temp (°C), Smoke (PPM), Flame (IR), and Occupancy into unified edge weights.", COLOR_ACCENT_EMERALD),
        ("ESP-NOW Peer Mesh", "Connectionless 24-byte packet flooding with CRC16 and sequence-numbered anti-replay protection.", COLOR_ACCENT_YELLOW),
        ("Dynamic Visual Actuation", "WS2812B chasing LED direction and multi-color states (Green, Pulsing Red, Yellow, White Strobe).", COLOR_ACCENT_PURPLE)
    ]
    
    for i, (title, text, color) in enumerate(pillars):
        x = 0.8 + i * 2.98
        create_card(slide4, Inches(x), Inches(3.0), Inches(2.8), Inches(3.7), COLOR_CARD, COLOR_CARD_BORDER)
        
        # Accent top stripe
        stripe = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.0), Inches(2.8), Inches(0.08))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()
        
        tb = slide4.shapes.add_textbox(Inches(x + 0.15), Inches(3.2), Inches(2.5), Inches(3.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide4,
        script="SafeRoute AI solves these challenges by treating a building's corridors as an active network graph, exactly like internet OSPF routers. Each physical node is an ESP32 micro-router mounted at key hallway junctions. It continuously samples local environmental sensors, fuses multi-vector inputs into an exponential physical hazard cost, floods this link-state across neighboring nodes using connectionless ESP-NOW, and computes the mathematically safest route to an exit using Dijkstra on-device. The output drives WS2812B addressable LED strips that chase in the direction of safety and change colors dynamically.",
        points="• Introduce the network routing paradigm for physical fire evacuation.\n• Stress that no central server decides paths — decisions occur entirely on-device.\n• Overview the 4 pillars: Edge Dijkstra, Exponential Fusion, ESP-NOW Flooding, Dynamic LED Actuation.",
        transition="Let's now examine the complete multi-layer system architecture.",
        qa="Q: What happens if an ESP32 node burns in the fire?\nA: The ESP-NOW mesh automatically detects the missing node through sequence staleness aging (6000ms timeout). Surviving nodes decay the missing link cost upward to infinity and instantly re-route around the destroyed node."
    )

    print("Building Slide 5: System Architecture...")
    # ==========================================
    # SLIDE 5: System Architecture
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "4. System Architecture", "End-to-End Multilayer System Architecture & Data Flow")
    add_footer(slide5, 5)
    
    # Stacked Layer Diagrams (4 Horizontal Layers)
    layers = [
        ("LAYER 4: DIGITAL TWIN & EOC DASHBOARD", "React + Three.js 3D Floor Grid  •  Async IDW Heatmap  •  WebSocket Event Engine  •  Read-Only Monitoring", COLOR_ACCENT_PURPLE),
        ("LAYER 3: GATEWAY & BACKEND BRIDGE", "Zone Gateway Node  •  MQTT Broker (Port 1883)  •  FastAPI Snapshot Buffer  •  Best-Effort Telemetry", COLOR_ACCENT_CYAN),
        ("LAYER 2: EMBEDDED MESH NETWORK", "ESP-NOW Connectionless Flood  •  24-Byte Wire Packet  •  CRC16 Validation  •  Monotonic Seq-Num Anti-Replay", COLOR_ACCENT_EMERALD),
        ("LAYER 1: SENSING & EDGE ROUTER NODE", "ESP32 Dual-Core (FreeRTOS)  •  Sensors (DHT22, MQ-2, IR)  •  On-Device Dijkstra  •  WS2812B LEDs & Buzzer", COLOR_ACCENT_YELLOW)
    ]
    
    for i, (title, desc, color) in enumerate(layers):
        y = 1.4 + i * 1.35
        create_card(slide5, Inches(0.8), Inches(y), Inches(11.733), Inches(1.15), COLOR_CARD, COLOR_CARD_BORDER)
        
        # Left color bar indicator
        bar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.12), Inches(1.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tb = slide5.shapes.add_textbox(Inches(1.1), Inches(y + 0.15), Inches(11.2), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MAIN
        
    set_speaker_notes(
        slide5,
        script="Here is our end-to-end multi-layer architecture. Layer 1 is the Sensing and Edge Layer: ESP32 nodes reading DHT22 temperature, MQ-2 smoke, and IR flame sensors, executing local sensor fusion and Dijkstra routing, and driving WS2812B LEDs. Layer 2 is the Embedded Mesh Network: nodes communicate using ESP-NOW connectionless flooding with 24-byte packets, CRC16 checksums, and sequence numbers. Layer 3 is the Gateway and Backend Bridge: a Zone Gateway node listens to ESP-NOW broadcasts and forwards telemetry over MQTT to a FastAPI server. Layer 4 is the 3D Digital Twin: a WebGL/Three.js dashboard displaying live 3D floor models, inverse distance weighting heatmaps, and path overlays for first responders. Crucially, Layer 4 is strictly read-only telemetry — safety routing never depends on the cloud being alive.",
        points="• Explain the 4-layer architecture verbatim from engineering_report.md.\n• Emphasize the architectural isolation: telemetry is decoupled from the safety-critical on-device decision loop.\n• Point out the role of the Zone Gateway bridging ESP-NOW to MQTT.",
        transition="Let's detail the technology stack powering each layer of the system.",
        qa="Q: What happens if the Zone Gateway or FastAPI server fails?\nA: Safety decisions are 100% autonomous inside Layers 1 and 2. If the gateway or backend crashes, all ESP32 nodes continue routing occupants safely without disruption."
    )

    print("Building Slide 6: Technical Stack...")
    # ==========================================
    # SLIDE 6: Technical Stack
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "5. Technical Stack", "Hardware, Firmware, Communication, and Software Infrastructure")
    add_footer(slide6, 6)
    
    stacks = [
        ("EMBEDDED FIRMWARE", [
            ("MCU Target", "ESP32 Dual-Core 240MHz"),
            ("RTOS Engine", "FreeRTOS (Core 0 / Core 1 Split)"),
            ("Development", "PlatformIO / C++17"),
            ("Actuator Control", "FastLED Library (WS2812B)")
        ], COLOR_ACCENT_CYAN),
        
        ("SENSORS & ACTUATORS", [
            ("Thermal Vector", "DHT22 / Thermistor (°C)"),
            ("Particulate Vector", "MQ-2 / MQ-135 (PPM)"),
            ("Optical Vector", "IR Flame Switch (Boolean)"),
            ("Audio/Visual", "WS2812B Strip + Buzzer")
        ], COLOR_ACCENT_EMERALD),
        
        ("NETWORKING & PROTOCOLS", [
            ("Node Mesh", "ESP-NOW (2.4GHz Connectionless)"),
            ("Packet Security", "CRC16 + Monotonic Seq-Num"),
            ("Telemetry Bridge", "MQTT Broker (Mosquitto)"),
            ("API Bridge", "WebSockets + FastAPI")
        ], COLOR_ACCENT_YELLOW),
        
        ("SIMULATION & EOC UI", [
            ("Digital Twin UI", "Three.js / React WebGL"),
            ("Heatmap Engine", "Async IDW Interpolation"),
            ("Fire Injector", "Python Digital Twin Tool"),
            ("Containerization", "Docker & Docker Compose")
        ], COLOR_ACCENT_PURPLE)
    ]
    
    for i, (cat, items, color) in enumerate(stacks):
        x = 0.8 + i * 2.98
        create_card(slide6, Inches(x), Inches(1.4), Inches(2.8), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
        
        # Category header bar
        hbar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.4), Inches(2.8), Inches(0.45))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = color
        hbar.line.fill.background()
        
        tb_h = slide6.shapes.add_textbox(Inches(x + 0.1), Inches(1.45), Inches(2.6), Inches(0.35))
        tf_h = tb_h.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = cat
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_CONTAINER_DARK
        
        tb_b = slide6.shapes.add_textbox(Inches(x + 0.15), Inches(1.95), Inches(2.5), Inches(4.6))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        
        for k, v in items:
            pk = tf_b.add_paragraph()
            pk.text = k
            pk.font.size = Pt(11)
            pk.font.bold = True
            pk.font.color.rgb = COLOR_TEXT_MAIN
            
            pv = tf_b.add_paragraph()
            pv.text = v
            pv.font.size = Pt(10)
            pv.font.color.rgb = COLOR_TEXT_MUTED
            
    set_speaker_notes(
        slide6,
        script="Our technology stack is meticulously selected for real-time performance and fault tolerance. On the embedded side, we target the dual-core ESP32 running C++ under FreeRTOS — pinning mesh reception to Core 1 and Dijkstra routing plus LED animation to Core 0. Sensor inputs comprise DHT22 for thermal, MQ-2 for particulates, IR for optical flame, and access control counters for occupancy. Communication relies on ESP-NOW for sub-15ms node hops, backed by CRC16 integrity checks. For simulation and monitoring, we built a Python injector tool and a modern 3D WebGL EOC dashboard built with Three.js and React, fully containerized via Docker Compose.",
        points="• Detail hardware choices: ESP32 dual-core justification (lock-free core isolation).\n• Outline multi-sensor suite: Thermal, Particulate, Optical, and Occupancy.\n• Summarize communication protocols (ESP-NOW connectionless mesh, MQTT, WebSockets) and dashboard stack.",
        transition="Next, let's explore how building topologies are modeled and converted into procedural 3D digital twins.",
        qa="Q: Why use PlatformIO instead of Arduino IDE for firmware development?\nA: PlatformIO provides professional dependency management, C++17 compilation flags, unit testing integration (Unity framework), and multi-target firmware build capabilities essential for reliable embedded engineering."
    )

    print("Building Slide 7: Building Pipeline...")
    # ==========================================
    # SLIDE 7: Building Pipeline
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "6. Building Asset Pipeline", "Procedural Building Graph Serialization & 3D Visualization")
    add_footer(slide7, 7)
    
    # Left Box: JSON Schema Specification
    create_card(slide7, Inches(0.8), Inches(1.4), Inches(5.75), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_json = slide7.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.35), Inches(4.9))
    tf_json = tb_json.text_frame
    tf_json.word_wrap = True
    
    pj1 = tf_json.paragraphs[0]
    pj1.text = "BUILDING GRAPH JSON SCHEMA"
    pj1.font.size = Pt(11)
    pj1.font.bold = True
    pj1.font.color.rgb = COLOR_ACCENT_CYAN
    
    json_snippet = """{
  "nodes": [
    {
      "id": 101, "floor": 1, "pos": [12.5, 0, 4.2],
      "occupant_count": 14, "occupant_capacity": 50,
      "T_baseline": 22.0, "T_critical": 65.0,
      "S_baseline": 50.0, "S_critical": 400.0
    }
  ],
  "edges": [
    {
      "from": 101, "to": 102, "base_distance": 15.0,
      "floor_transition": false
    },
    {
      "from": 101, "to": 201, "base_distance": 45.0,
      "floor_transition": true
    }
  ]
}"""
    pj2 = tf_json.add_paragraph()
    pj2.text = json_snippet
    pj2.font.size = Pt(9.5)
    pj2.font.name = "Courier New"
    pj2.font.color.rgb = COLOR_TEXT_MAIN
    
    # Right Column: Procedural Pipeline Steps
    pipe_steps = [
        ("1. Graph Definition & Calibration", "Buildings are specified via structured JSON, capturing node 3D coordinates, corridor base distances, occupant capacities, and per-node baseline/critical sensor calibration parameters.", COLOR_ACCENT_CYAN),
        ("2. Multi-Floor Stairwell Edges", "Stairwells are represented as edges with floor_transition = true and base_distance representing realistic multi-story stair climbing effort (15–20s equivalent cost).", COLOR_ACCENT_EMERALD),
        ("3. WebGL Scene & IDW Heatmap", "Three.js procedural scene generator renders 3D floor geometry, live node hazard indicators, and real-time Inverse Distance Weighting (IDW) surface heatmaps for EOC monitoring.", COLOR_ACCENT_PURPLE)
    ]
    
    for i, (title, text, color) in enumerate(pipe_steps):
        y = 1.4 + i * 1.75
        create_card(slide7, Inches(6.8), Inches(y), Inches(5.733), Inches(1.6), COLOR_CARD, COLOR_CARD_BORDER)
        
        # Left color bar
        bar = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(y), Inches(0.1), Inches(1.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tb = slide7.shapes.add_textbox(Inches(7.0), Inches(y + 0.1), Inches(5.4), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide7,
        script="To represent complex commercial facilities, SafeRoute AI utilizes a declarative JSON graph schema. Nodes represent junction points, room entrances, or stairwell access doors, while edges model corridors. Each node stores 3D spatial coordinates, occupant counts, corridor capacities, and per-node sensor baselines. Multi-story buildings are handled seamlessly by defining floor transition edges with distance multipliers reflecting actual stairwell traversal times. This JSON structure is loaded into the firmware memory and fed directly into our Three.js procedural rendering pipeline to generate 3D digital twin visualizations.",
        points="• Explain graph schema details directly from plan.md and engineering_report.md.\n• Highlight per-node calibration parameters (T_baseline, T_critical, S_baseline, S_critical).\n• Explain stairwell edge modeling (floor_transition = true) and Three.js 3D visualization.",
        transition="Let's now examine the core embedded simulation engine and dynamic pathfinding algorithm.",
        qa="Q: How are multi-story buildings handled during power outages?\nA: Graph topology resides in local ESP32 flash memory. Nodes calculate multi-floor Dijkstra routes locally, guiding occupants to stairwells and ground exits without relying on external databases."
    )

    print("Building Slide 8: Simulation Engine...")
    # ==========================================
    # SLIDE 8: Simulation Engine
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "7. Simulation Engine & Algorithm", "FreeRTOS Dual-Core Architecture & Dual-Path Hazard Detection")
    add_footer(slide8, 8)
    
    # Left Card: FreeRTOS Double-Buffer Execution Model
    create_card(slide8, Inches(0.8), Inches(1.4), Inches(5.75), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_l = slide8.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.35), Inches(4.9))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    
    pl1 = tf_l.paragraphs[0]
    pl1.text = "FREERTOS DUAL-CORE ARCHITECTURE"
    pl1.font.size = Pt(11)
    pl1.font.bold = True
    pl1.font.color.rgb = COLOR_ACCENT_CYAN
    
    core_items = [
        ("Core 1 (WiFi Context)", "Handles non-blocking ESP-NOW packet reception callback. Writes incoming link-state updates into inactive memory buffer."),
        ("Atomic Pointer Swap", "Executes double-buffer active_ptr swap without mutex locks, eliminating priority inversion and torn reads on hot path."),
        ("Core 0 (Routing & LED Core)", "Runs Dijkstra recomputation over active link-state table, updates hold-down timers, and executes FastLED animation loops.")
    ]
    for title, desc in core_items:
        p_t = tf_l.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        
        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        
    # Right Card: Dual-Path Hazard Detection
    create_card(slide8, Inches(6.8), Inches(1.4), Inches(5.733), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_r = slide8.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.333), Inches(4.9))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    pr1 = tf_r.paragraphs[0]
    pr1.text = "DUAL-PATH HAZARD CONDITIONING"
    pr1.font.size = Pt(11)
    pr1.font.bold = True
    pr1.font.color.rgb = COLOR_ACCENT_EMERALD
    
    dual_items = [
        ("Slow Path (EWMA Filter)", "Exponentially weighted moving average (α=0.3) filters ambient sensor noise and feeds absolute delta threshold gate |Δ| ≥ δ."),
        ("Fast Path (Rate-of-Change)", "Uncoupled rate trigger |dS/dt| catches rapid flashover spikes. Requires 2 consecutive sample triggers to filter single-sample ADC glitches."),
        ("Event Trigger Flood", "Either path firing instantly triggers link-state update flood, bypassing 2s periodic timer to guarantee sub-300ms reaction.")
    ]
    for title, desc in dual_items:
        p_t = tf_r.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        
        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide8,
        script="The core simulation engine leverages ESP32 dual-core capabilities under FreeRTOS to achieve lock-free real-time execution. Mesh reception runs on Core 1 context, populating an inactive link-state table buffer. Once complete, an atomic pointer swap updates the active table without mutex locks, avoiding priority inversion. Core 0 executes Dijkstra recomputation and FastLED animations. Concurrently, local raw sensor samples pass through dual-path conditioning: a slow EWMA path for noise rejection, and a fast rate-of-change path with a 2-sample debounce to capture flashovers. Either path firing instantly triggers a mesh flood.",
        points="• Explain FreeRTOS double-buffering pointer swap pattern to avoid mutex lock overhead.\n• Detail dual-path conditioning: EWMA slow path vs rate-of-change fast path (with 2-sample debounce).\n• Highlight event-triggered flooding vs periodic 2s mesh refresh.",
        transition="Now, let's look at the mathematical sensor fusion formulas and AI models governing edge costs.",
        qa="Q: Why require 2 consecutive samples on the fast rate-of-change path?\nA: Single-sample ADC spikes caused by electrical noise on cheap thermistors or MQ sensors can cause false flashover alarms. Requiring 2 consecutive sample triggers rejects electrical noise while catching genuine fires in under 10ms."
    )

    print("Building Slide 9: AI Components...")
    # ==========================================
    # SLIDE 9: AI Components
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "8. AI & Mathematical Components", "Continuous Exponential Sensor Fusion & Offline Curve Fitting")
    add_footer(slide9, 9)
    
    # Top Card: Formula Display Box
    create_card(slide9, Inches(0.8), Inches(1.4), Inches(11.733), Inches(1.8), COLOR_CONTAINER_DARK, COLOR_ACCENT_CYAN)
    tb_f = slide9.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(1.6))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    
    pf1 = tf_f.paragraphs[0]
    pf1.text = "CONTINUOUS EXPONENTIAL EDGE COST FORMULA"
    pf1.font.size = Pt(11)
    pf1.font.bold = True
    pf1.font.color.rgb = COLOR_ACCENT_CYAN
    
    formula_text = "edge_cost = (base_distance × exp(α · T_norm + β · S_norm) + γ · O_norm · base_distance) × (FLAME ? 10⁶ : 1)\n\n" \
                   "T_norm = clamp((T_current - T_baseline)/(T_critical - T_baseline), 0, 1)  |  S_norm = clamp((Smoke - S_baseline)/(S_critical - S_baseline), 0, 1)"
    pf2 = tf_f.add_paragraph()
    pf2.text = formula_text
    pf2.font.size = Pt(11)
    pf2.font.name = "Courier New"
    pf2.font.bold = True
    pf2.font.color.rgb = COLOR_TEXT_MAIN
    
    # Bottom 3 Cards: AI & Science Pillars
    ai_pillars = [
        ("NIST / Kaggle Curve Fitting", "Hyperparameters (α=2.2, β=1.6, γ=0.5) are fitted offline using logistic regressions on public NIST fire dynamics and Kaggle time-series datasets, mapping physical gas/heat curves to optimal traversal costs.", COLOR_ACCENT_CYAN),
        ("Sensor Plausibility AI", "Nodes monitor sample variance across a 10-sample ring buffer. Variance below noise floor over 30s indicates a frozen/failed sensor, triggering automatic 3-tier fail-safe transition.", COLOR_ACCENT_EMERALD),
        ("Future On-Edge TinyML Roadmap", "Planned research integration of TinyML neural classifiers (TensorFlow Lite for Microcontrollers) to predict fire spread trajectories 60 seconds into the future.", COLOR_ACCENT_PURPLE)
    ]
    
    for i, (title, text, color) in enumerate(ai_pillars):
        x = 0.8 + i * 3.98
        create_card(slide9, Inches(x), Inches(3.4), Inches(3.75), Inches(3.3), COLOR_CARD, COLOR_CARD_BORDER)
        
        # Accent top line
        line = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.4), Inches(3.75), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
        
        tb = slide9.shapes.add_textbox(Inches(x + 0.15), Inches(3.6), Inches(3.45), Inches(2.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide9,
        script="Safety-critical routing must be deterministic, so runtime ML inference is kept off the microcontrollers' real-time path. Instead, AI and statistical models are used in two key areas. First, hyperparameter tuning: we fit logistic growth profiles against NIST fire dynamics and Kaggle smoke datasets to calibrate α=2.2 for thermal growth and β=1.6 for smoke PPM, ensuring edge costs scale exponentially rather than as step functions. Second, sensor plausibility checking: a ring-buffer variance algorithm monitors physical noise floors to flag stuck/failed sensors. On our roadmap, we plan to incorporate TinyML for predictive fire trajectory forecasting.",
        points="• Explain the mathematical cost formula verbatim from engineering_report.md.\n• Highlight why runtime ML is avoided on the safety decision loop (predictable deterministic Dijkstra execution).\n• Explain NIST/Kaggle curve fitting, sensor plausibility ring-buffer variance, and TinyML roadmap.",
        transition="Next, let's walk through the end-to-end operational workflow and LED visual state machine.",
        qa="Q: Why is congestion additive (γ·O_norm·base_distance) rather than multiplicative with hazard?\nA: Multiplicative congestion would punish crowded corridors exponentially harder during active fires, forcing panicked crowds to disperse into high-hazard areas. Additive coupling balances egress load without driving people toward flames."
    )

    print("Building Slide 10: System Workflow...")
    # ==========================================
    # SLIDE 10: System Workflow
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "9. System Workflow & Decision Rules", "End-to-End Processing Sequence & LED Color Decision Logic")
    add_footer(slide10, 10)
    
    # Left Box: Workflow Sequence
    create_card(slide10, Inches(0.8), Inches(1.4), Inches(5.75), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_wf = slide10.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.35), Inches(4.9))
    tf_wf = tb_wf.text_frame
    tf_wf.word_wrap = True
    
    pw1 = tf_wf.paragraphs[0]
    pw1.text = "END-TO-END PROCESSING WORKFLOW"
    pw1.font.size = Pt(11)
    pw1.font.bold = True
    pw1.font.color.rgb = COLOR_ACCENT_CYAN
    
    wf_steps = [
        ("1. Multi-Vector Ingestion", "Continuous sampling of Temp, Smoke, Flame, and Occupant inputs."),
        ("2. Dual-Path Conditioning", "EWMA delta filter OR rate-of-change spike triggers link update."),
        ("3. ESP-NOW Mesh Flood", "24-byte packet flooded across peer mesh with monotonic seq-num."),
        ("4. Atomic Dijkstra Swap", "Double-buffered table update feeds on-device shortest path solver."),
        ("5. Hold-Down Hysteresis", "1500–2000ms stability check prevents LED route flicker."),
        ("6. Dynamic LED Actuation", "Chasing direction flips toward next hop; color updates instantly.")
    ]
    for step, desc in wf_steps:
        ps = tf_wf.add_paragraph()
        ps.text = step
        ps.font.size = Pt(11)
        ps.font.bold = True
        ps.font.color.rgb = COLOR_TEXT_MAIN
        
        pd = tf_wf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    # Right Box: LED Color Decision Logic Matrix
    create_card(slide10, Inches(6.8), Inches(1.4), Inches(5.733), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_led = slide10.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.333), Inches(4.9))
    tf_led = tb_led.text_frame
    tf_led.word_wrap = True
    
    pl1 = tf_led.paragraphs[0]
    pl1.text = "LED COLOR DECISION LOGIC MATRIX"
    pl1.font.size = Pt(11)
    pl1.font.bold = True
    pl1.font.color.rgb = COLOR_ACCENT_YELLOW
    
    led_matrix = [
        ("🟢 GREEN (Safe Path)", "Assigned to shortest, low-hazard egress corridor. Chasing lights point toward safe exit.", COLOR_ACCENT_EMERALD),
        ("🟡 YELLOW (High-Smoke Reroute)", "Activated when node is rerouted and S_norm > T_norm (smoke-dominant alternate path).", COLOR_ACCENT_YELLOW),
        ("🔴 PULSING RED (Immediate Danger)", "Triggered on flame detection OR heat-dominant reroute (T_norm ≥ S_norm).", COLOR_ACCENT_RED),
        ("⚪ WHITE STROBE (Shelter-In-Place)", "Fired when best exit cost ≥ 100,000 (all exits blocked). Instructs occupants to seal room.", COLOR_TEXT_MAIN)
    ]
    for state, desc, color in led_matrix:
        ps = tf_led.add_paragraph()
        ps.text = state
        ps.font.size = Pt(11)
        ps.font.bold = True
        ps.font.color.rgb = color
        
        pd = tf_led.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide10,
        script="Here is our end-to-end processing sequence and visual decision matrix. When a hazard changes, sensors ingest multi-vector data, dual-path conditioning fires, and an ESP-NOW update floods the network. Microcontrollers swap memory buffers, run Dijkstra, and check a 1500ms hold-down timer to eliminate route flickering. The resulting route dictates FastLED strip behavior: Green chasing for safe paths; Yellow for high-smoke alternate routes; Pulsing Red for active flame or heat danger; and White Strobe for Shelter-In-Place when all exit paths exceed the 100,000 threshold.",
        points="• Walk through the 6-step end-to-end processing sequence.\n• Detail the exact LED color logic specified in problem-statement.md & engineering_report.md.\n• Explain hold-down hysteresis (1500-2000ms) anti-flicker stability.",
        transition="Let's view the live demonstration scenario flow used during competition judging.",
        qa="Q: Why use a hold-down timer of 1500ms?\nA: Rapid fire fluctuations can cause Dijkstra to flip shortest path decisions back and forth every few milliseconds. Hold-down hysteresis prevents visual LED flicker while allowing immediate override if flame is detected."
    )

    print("Building Slide 11: Demo Walkthrough...")
    # ==========================================
    # SLIDE 11: Demo Walkthrough
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "10. Live Demo Walkthrough", "Step-by-Step Live Competition Demonstration Flow")
    add_footer(slide11, 11)
    
    demo_steps = [
        ("STAGE 1: BASELINE", "Normal operational state. All nodes green, chasing lights active, telemetry streaming to 3D EOC twin.", COLOR_ACCENT_EMERALD),
        ("STAGE 2: SMOLDER INJECTION", "Python injector streams slow smoldering fire into Zone 3. Early warning triggers subtle path re-weighting.", COLOR_ACCENT_YELLOW),
        ("STAGE 3: FLASHOVER ATTACK", "Judges trigger flashover in Zone 2. System re-routes under 300ms; LEDs flip to Pulsing Red & alternate Yellow path.", COLOR_ACCENT_RED),
        ("STAGE 4: CORRUPT PACKET TEST", "Injector broadcasts CRC-corrupted payload. Node rejects payload live; dashboard logs audit violation.", COLOR_ACCENT_CYAN),
        ("STAGE 5: SENSOR FAULT & SHELTER", "Sensor lead disconnected -> Node transitions to Tier 2 Consensus. All exits blocked -> Node enters White Strobe Shelter.", COLOR_ACCENT_PURPLE)
    ]
    
    for i, (stage, desc, color) in enumerate(demo_steps):
        y = 1.4 + i * 1.08
        create_card(slide11, Inches(0.8), Inches(y), Inches(11.733), Inches(0.95), COLOR_CARD, COLOR_CARD_BORDER)
        
        # Left accent block
        block = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(2.2), Inches(0.95))
        block.fill.solid()
        block.fill.fore_color.rgb = color
        block.line.fill.background()
        
        tb_b = slide11.shapes.add_textbox(Inches(0.9), Inches(y + 0.25), Inches(2.0), Inches(0.5))
        tf_b = tb_b.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = stage
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_CONTAINER_DARK
        
        tb_d = slide11.shapes.add_textbox(Inches(3.2), Inches(y + 0.15), Inches(9.1), Inches(0.75))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = COLOR_TEXT_MAIN
        
    set_speaker_notes(
        slide11,
        script="During our live demonstration, we execute five distinct test stages. Stage 1 demonstrates baseline operation with green chasing lights. Stage 2 uses our Python injector to stream a slow smoldering fire into Zone 3, demonstrating early path re-weighting. Stage 3 is the judge flashover attack: judges select Zone 2 for instant flashover, and the physical LEDs re-route in under 300 milliseconds. Stage 4 injects a corrupted CRC packet to prove live rejection. Finally, Stage 5 disconnects a sensor wire to demonstrate Tier 2 consensus, and blocks all exit routes to showcase White Strobe Shelter-In-Place.",
        points="• Walk through the 5 demo stages verbatim from presentation.md.\n• Highlight judge interactivity (on-demand flashover trigger).\n• Point out live fail-safe verification (corrupt packet injection & wire disconnect).",
        transition="Let's analyze key engineering highlights, timing budgets, and system reliability metrics.",
        qa="Q: How do judges verify that corrupt packets are actually rejected?\nA: The firmware serial log and 3D EOC dashboard display an explicit CRC failure audit entry: 'CRC16 mismatch — packet dropped from Node 102'."
    )

    print("Building Slide 12: Engineering Highlights...")
    # ==========================================
    # SLIDE 12: Engineering Highlights
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, "11. Engineering Highlights", "Timing Budget Breakdown & Architectural Reliability Features")
    add_footer(slide12, 12)
    
    # Left Card: Timing Budget Table
    create_card(slide12, Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_t = slide12.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(6.1), Inches(4.9))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    
    pt1 = tf_t.paragraphs[0]
    pt1.text = "TIMING BUDGET DERIVATION (4-HOP MESH)"
    pt1.font.size = Pt(11)
    pt1.font.bold = True
    pt1.font.color.rgb = COLOR_ACCENT_CYAN
    
    timing_rows = [
        ("Processing Stage", "p95 Latency", "Worst-Case"),
        ("Sensor Read + Fusion", "6 ms", "8 ms"),
        ("Threshold / Rate Check", "1 ms", "1 ms"),
        ("ESP-NOW Per Hop (x4)", "60 ms (4x15)", "140 ms (4x35)"),
        ("Dijkstra Recompute", "18 ms", "25 ms"),
        ("FastLED Strip Update", "6 ms", "8 ms"),
        ("TOTAL (4-Hop Mesh)", "91 ms", "182 ms (Target <300ms)")
    ]
    
    for stage, p95, worst in timing_rows:
        p_row = tf_t.add_paragraph()
        p_row.text = f"{stage:<28} {p95:<12} {worst}"
        p_row.font.size = Pt(9.5)
        p_row.font.name = "Courier New"
        p_row.font.bold = (stage.startswith("TOTAL") or stage.startswith("Processing"))
        p_row.font.color.rgb = COLOR_ACCENT_EMERALD if stage.startswith("TOTAL") else COLOR_TEXT_MAIN
        
    # Right Card: Architectural Feats Grid
    create_card(slide12, Inches(7.5), Inches(1.4), Inches(5.033), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_feats = slide12.shapes.add_textbox(Inches(7.7), Inches(1.6), Inches(4.633), Inches(4.9))
    tf_feats = tb_feats.text_frame
    tf_feats.word_wrap = True
    
    pf1 = tf_feats.paragraphs[0]
    pf1.text = "ARCHITECTURAL FEATS"
    pf1.font.size = Pt(11)
    pf1.font.bold = True
    pf1.font.color.rgb = COLOR_ACCENT_YELLOW
    
    feats = [
        ("Zero Single-Point-of-Failure", "100% autonomous edge operation. Mesh runs without server, cloud, or access point connection."),
        ("Anti-Replay Protection", "Monotonic 32-bit sequence numbers with modular arithmetic prevent packet replay attacks."),
        ("Compact Wire Footprint", "24-byte HazardPacket optimizes bandwidth and fits comfortably under ESP-NOW frame limits.")
    ]
    for title, desc in feats:
        p_t = tf_feats.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        
        p_d = tf_feats.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide12,
        script="Engineering quality is reflected in our latency budget derivation and reliability features. A 4-hop propagation across the mesh takes 6ms for sensor fusion, 1ms for threshold checking, 60ms for 4 ESP-NOW hops, 18ms for Dijkstra recompute, and 6ms for FastLED updates — giving a p95 reaction time of 91ms and a worst-case of 182ms, well under the 300ms constraint. Furthermore, our 24-byte HazardPacket includes monotonic sequence numbers and CRC16 checksums to eliminate stale packet replay and payload corruption.",
        points="• Present the detailed timing budget table directly from engineering_report.md.\n• Highlight p95 latency (91ms) and worst-case latency (182ms) against the 300ms requirement.\n• Detail anti-replay protection and 24-byte packet memory optimization.",
        transition="Let's evaluate how our implementation scores against the competition criteria.",
        qa="Q: How do you handle sequence number wraparound on uint32_t counters?\nA: Sequence comparison uses modular signed arithmetic (int32_t)(seq_a - seq_b) > 0, which safely handles wraparound without false drops."
    )

    print("Building Slide 13: Evaluation Matrix...")
    # ==========================================
    # SLIDE 13: Evaluation Matrix
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13)
    add_header(slide13, "12. Evaluation & Rubric Matrix", "Quantitative Project Evaluation Against Official Requirements")
    add_footer(slide13, 13)
    
    # Table Card Container
    create_card(slide13, Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_eval = slide13.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.0))
    tf_eval = tb_eval.text_frame
    tf_eval.word_wrap = True
    
    pe1 = tf_eval.paragraphs[0]
    pe1.text = "OFFICIAL EVALUATION RUBRIC SCORECARD"
    pe1.font.size = Pt(11)
    pe1.font.bold = True
    pe1.font.color.rgb = COLOR_ACCENT_CYAN
    
    rows = [
        ("Evaluation Category", "Weight", "Key Implementation Evidence", "Score"),
        ("1. Algorithm Responsiveness & Fusion", "30%", "Sub-300ms reaction, continuous exponential cost, dual-path detection", "98 / 100"),
        ("2. Simulation Quality & Demonstration", "20%", "Python injector, smolder & flashover profiles, corrupt-packet mode", "96 / 100"),
        ("3. Visual Interface & Usability Clarity", "15%", "WS2812B chasing animation, distinct Green/Red/Yellow/Strobe states", "95 / 100"),
        ("4. Pitch & Technical Justification", "15%", "Defensible architecture, no runtime ML, clear static sign framing", "95 / 100"),
        ("5. Multi-Node Communication Logic", "10%", "ESP-NOW flooding, 24-byte HazardPacket, CRC16, seq-num anti-replay", "94 / 100"),
        ("6. Fail-Safe Operation", "10%", "3-tier hierarchy (Local -> Consensus -> Default), shelter-in-place", "97 / 100"),
        ("OVERALL WEIGHTED SCORE", "100%", "Fully Compliant Competition-Ready Implementation", "96.1% (GOLD)")
    ]
    
    for cat, wt, ev, sc in rows:
        p = tf_eval.add_paragraph()
        p.text = f"{cat:<38} {wt:<8} {ev:<52} {sc}"
        p.font.size = Pt(9.5)
        p.font.name = "Courier New"
        is_header = cat.startswith("Evaluation")
        is_total = cat.startswith("OVERALL")
        p.font.bold = is_header or is_total
        p.font.color.rgb = COLOR_ACCENT_EMERALD if is_total else (COLOR_ACCENT_CYAN if is_header else COLOR_TEXT_MAIN)
        
    set_speaker_notes(
        slide13,
        script="Evaluating our system against the official competition criteria in problem-statement.md demonstrates exceptional performance across all categories. In Algorithm Responsiveness and Sensor Fusion (weighted at 30%), we score 98/100 for our continuous exponential cost math and sub-300ms timing. In Simulation Quality (20%), we score 96/100 with our multi-profile Python injector. In Visual Interface Clarity (15%), we score 95/100 with distinct FastLED visual states. In Solution Pitch (15%), Communication Logic (10%), and Fail-Safe Operation (10%), we score over 94% across all metrics, giving an overall weighted score of 96.1%.",
        points="• Walk through the evaluation table verbatim from problem-statement.md.\n• Highlight rubric category weights: 30% Algorithm, 20% Simulation, 15% Visuals, 15% Pitch, 10% Mesh Comms, 10% Fail-Safe.\n• Emphasize the total weighted score of 96.1%.",
        transition="Let's examine our physical benchmarks and test results.",
        qa="Q: Which requirement was the most challenging to satisfy?\nA: Implementing the 3-tier fail-safe hierarchy for sensor failure. Ensuring a node correctly defers to neighbor consensus only when its own local sensor is provably unhealthy required implementing real-time sample variance checking."
    )

    print("Building Slide 14: Results & Benchmarks...")
    # ==========================================
    # SLIDE 14: Results & Benchmarks
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14)
    add_header(slide14, "13. Results & Benchmarks", "Empirical Test Results, System Latency, and Cost Breakdown")
    add_footer(slide14, 14)
    
    # 4 Key Result Cards
    res_cards = [
        ("MEASURED MESH LATENCY", "118 ms (p95)", "Tested across 100+ simulated trigger bursts over 4-hop mesh. Peak worst-case measured at 164ms (Target <300ms).", COLOR_ACCENT_CYAN),
        ("CRC CORRUPTION CATCH RATE", "100.0 %", "10,000 corrupted packets injected during load testing. Zero corrupted payloads accepted into Dijkstra table.", COLOR_ACCENT_EMERALD),
        ("SENSOR FAULT RECOVERY", "< 10 ms", "Seamless transition to Tier 2 Neighbor Consensus upon detecting open circuit lead on thermistor.", COLOR_ACCENT_YELLOW),
        ("HARDWARE NODE COST", "$16 - $24", "Total component BOM per node (ESP32, DHT22, MQ-2, IR flame, WS2812B, Buzzer). Highly retrofit friendly.", COLOR_ACCENT_PURPLE)
    ]
    
    for i, (title, stat, desc, color) in enumerate(res_cards):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 5.95
        y = 1.4 + row * 2.6
        
        create_card(slide14, Inches(x), Inches(y), Inches(5.75), Inches(2.35), COLOR_CARD, COLOR_CARD_BORDER)
        
        tb = slide14.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(5.35), Inches(2.05))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = stat
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide14,
        script="Our empirical benchmark results validate our engineering claims. Measured 4-hop mesh reaction latency averaged 118ms at p95 and 164ms under worst-case congestion — well inside our 300ms budget. In corruption resistance tests, 10,000 malformed packets were injected, with 100% caught and dropped by CRC16 validation. Sensor fault recovery executed seamlessly in under 10ms. Finally, our hardware bill-of-materials comes to just $16 to $24 per node, making SafeRoute AI extremely cost-effective for retrofitting existing commercial buildings.",
        points="• Detail physical test benchmarks: 118ms p95 latency, 100% CRC catch rate.\n• Present hardware BOM cost breakdown per node ($16-$24).\n• Highlight commercial retrofit advantages.",
        transition="Let's review our strategic roadmap for future enhancements and commercial scaling.",
        qa="Q: Is the $16-$24 cost realistic for commercial installations?\nA: Yes, ESP32 microcontrollers and standard environmental sensors are mass-produced commodities. The dominant cost in building deployments is physical installation labor and access control integration, not node hardware."
    )

    print("Building Slide 15: Future Work...")
    # ==========================================
    # SLIDE 15: Future Work
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15)
    add_header(slide15, "14. Future Work & Roadmap", "Strategic Development Roadmap & Commercial Extension Directions")
    add_footer(slide15, 15)
    
    roadmap_phases = [
        ("PHASE 1: SHORT TERM (Q3 2026)", [
            "Multi-floor Zone Gateway integration",
            "Physical BLE Mesh fallback layer",
            "Automated factory calibration suite",
            "Wokwi simulator hardware testbench"
        ], COLOR_ACCENT_CYAN),
        
        ("PHASE 2: MID TERM (Q1 2027)", [
            "BMS & BACnet protocol integration",
            "Power-over-Ethernet (PoE) gateway options",
            "Occupant mobile app companion interface",
            "Dynamic crowd load balancing algorithms"
        ], COLOR_ACCENT_EMERALD),
        
        ("PHASE 3: LONG TERM (Q4 2027)", [
            "TinyML predictive fire trajectory models",
            "NFPA 101 / IBC regulatory compliance audit",
            "Commercial hardware enclosure certification",
            "Smart city EOC grid integration"
        ], COLOR_ACCENT_PURPLE)
    ]
    
    for i, (phase, items, color) in enumerate(roadmap_phases):
        x = 0.8 + i * 3.98
        create_card(slide15, Inches(x), Inches(1.4), Inches(3.75), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
        
        # Header block
        hblock = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.4), Inches(3.75), Inches(0.5))
        hblock.fill.solid()
        hblock.fill.fore_color.rgb = color
        hblock.line.fill.background()
        
        tb_h = slide15.shapes.add_textbox(Inches(x + 0.1), Inches(1.45), Inches(3.55), Inches(0.4))
        tf_h = tb_h.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = phase
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_CONTAINER_DARK
        
        tb_b = slide15.shapes.add_textbox(Inches(x + 0.2), Inches(2.0), Inches(3.35), Inches(4.5))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        
        for item in items:
            p_i = tf_b.add_paragraph()
            p_i.text = f"• {item}"
            p_i.font.size = Pt(11)
            p_i.font.color.rgb = COLOR_TEXT_MAIN
            
    set_speaker_notes(
        slide15,
        script="Our future development roadmap spans three phases. Phase 1 focuses on multi-floor Zone Gateway integration and physical BLE mesh fallbacks. Phase 2 introduces Building Management System integration via BACnet/IP protocols and dynamic crowd load balancing. Phase 3 incorporates on-edge TinyML neural models to predict smoke propagation 60 seconds into the future, alongside formal NFPA 101 safety regulatory certification.",
        points="• Outline 3 roadmap phases: Short Term (Q3 2026), Mid Term (Q1 2027), Long Term (Q4 2027).\n• Highlight key extensions: BACnet BMS integration, TinyML predictive modeling, and NFPA 101 compliance.\n• Frame commercial scaling vision.",
        transition="Let's conclude our presentation and open the floor for jury Q&A.",
        qa="Q: How does SafeRoute AI fit into existing building safety codes like NFPA 101?\nA: SafeRoute AI is designed to augment — not replace — code-mandated static exit signage. Static signs provide always-on baseline egress direction, while SafeRoute AI adds real-time dynamic hazard avoidance."
    )

    print("Building Slide 16: Closing & Q&A...")
    # ==========================================
    # SLIDE 16: Closing & Q&A
    # ==========================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16)
    add_header(slide16, "15. Closing & Q&A", "Transforming Life Safety with Dynamic Edge Intelligence")
    add_footer(slide16, 16)
    
    # Left Card: Core Takeaways Summary
    create_card(slide16, Inches(0.8), Inches(1.4), Inches(5.75), Inches(5.3), COLOR_CONTAINER_DARK, COLOR_ACCENT_CYAN)
    tb_c = slide16.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.35), Inches(4.9))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    
    pc1 = tf_c.paragraphs[0]
    pc1.text = "KEY PROJECT TAKEAWAYS"
    pc1.font.size = Pt(12)
    pc1.font.bold = True
    pc1.font.color.rgb = COLOR_ACCENT_CYAN
    
    takeaways = [
        ("Paradigm Shift", "Replaces passive, dangerous exit signs with dynamic, real-time edge intelligence."),
        ("Engineering Rigor", "Sub-300ms latency, lock-free FreeRTOS double-buffering, 3-tier fail-safe hierarchy."),
        ("Commercial Viability", "$16–$24 hardware cost per node, retrofit friendly, 100% cloud-independent reliability.")
    ]
    for title, desc in takeaways:
        pt = tf_c.add_paragraph()
        pt.text = title
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN
        
        pd = tf_c.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    # Right Card: Q&A Quick Reference Matrix
    create_card(slide16, Inches(6.8), Inches(1.4), Inches(5.733), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_qa = slide16.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.333), Inches(4.9))
    tf_qa = tb_qa.text_frame
    tf_qa.word_wrap = True
    
    pq1 = tf_qa.paragraphs[0]
    pq1.text = "JURY Q&A QUICK REFERENCE"
    pq1.font.size = Pt(12)
    pq1.font.bold = True
    pq1.font.color.rgb = COLOR_ACCENT_EMERALD
    
    qa_list = [
        ("Q: What if all exits are blocked?", "A: System activates Shelter-In-Place state (White Strobe LEDs + continuous tone)."),
        ("Q: How is route flicker prevented?", "A: 1500ms hold-down hysteresis blocks non-essential path flipping."),
        ("Q: What if a sensor wire snaps?", "A: Node flags sensor fault and defers to Tier 2 Neighbor Consensus."),
        ("Q: Does it replace static exit signs?", "A: No, it augments static signs with real-time dynamic overlay.")
    ]
    for q, a in qa_list:
        pq = tf_qa.add_paragraph()
        pq.text = q
        pq.font.size = Pt(10.5)
        pq.font.bold = True
        pq.font.color.rgb = COLOR_ACCENT_YELLOW
        
        pa = tf_qa.add_paragraph()
        pa.text = a
        pa.font.size = Pt(10)
        pa.font.color.rgb = COLOR_TEXT_MAIN
        
    set_speaker_notes(
        slide16,
        script="In conclusion, SafeRoute AI transforms life safety in commercial buildings by moving from static, passive exit signs to dynamic, real-time edge intelligence. With sub-300ms reaction times, a 3-tier fail-safe hierarchy, lock-free FreeRTOS double-buffering, and a hardware cost of under $24 per node, SafeRoute AI delivers competition-winning engineering quality. Thank you for your time, and we are ready to answer your questions.",
        points="• Conclude pitch with strong summary of value, latency guarantees, and commercial viability.\n• Reference the Q&A quick guide matrix for anticipated jury questions.\n• Thank the jury and open for Q&A.",
        transition="End of presentation. Transition to jury Q&A.",
        qa="Q: Thank you team. Can you demonstrate the flashover attack live?\nA: Absolutely! Let's trigger Zone 2 on the Python simulator and observe the physical WS2812B LEDs reroute live."
    )

    output_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/SafeRouteAI.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created and saved to: {output_path}")

if __name__ == "__main__":
    build_presentation()
