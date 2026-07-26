#!/usr/bin/env python3
import sys
from pptx import Presentation

def export_executive_notes():
    pptx_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/SafeRouteAI_Executive_6Slides.pptx"
    md_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/docs/executive_speaker_notes.md"
    
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    
    md_content = []
    md_content.append("# SafeRoute AI — Executive Pitch Deck Speaker Notes (6-Slide Edition)\n")
    md_content.append("**Project**: SafeRoute AI — Decentralized Evacuation Routing with Real-Time Hazard Mapping")
    md_content.append("**Target Pitch Time**: 3–5 Minutes (Fast-Paced Competition Pitch)")
    md_content.append(f"**Total Slides**: {total_slides} Widescreen (16:9) Dark-Theme Slides\n")
    md_content.append("---\n")
    
    for i, slide in enumerate(prs.slides, 1):
        title = f"Slide {i}"
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                text = shape.text_frame.text.split('\n')[0]
                if len(text) > 3 and not text.startswith("SafeRoute AI") and not text.startswith("16:9"):
                    title = f"Slide {i}: {text}"
                    break
                    
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text
            
        md_content.append(f"## {title}\n")
        if notes_text:
            md_content.append(f"```text\n{notes_text}\n```\n")
        else:
            md_content.append("> [!WARNING]\n> No speaker notes found on this slide!\n")
        md_content.append("---\n")
        
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md_content))
        
    print(f"Successfully exported executive speaker notes to: {md_path}")

if __name__ == "__main__":
    export_executive_notes()
