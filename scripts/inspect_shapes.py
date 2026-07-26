#!/usr/bin/env python3
import sys
from pptx import Presentation

def inspect_shapes():
    pptx_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/IDEA_Presentation_Format.pptx"
    prs = Presentation(pptx_path)
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"=== SLIDE {i} Shapes ===")
        for j, shape in enumerate(slide.shapes):
            text = shape.text_frame.text if shape.has_text_frame else "[No Text]"
            print(f"  Shape {j}: name='{shape.name}', type={shape.shape_type}, left={shape.left.inches:.2f}\", top={shape.top.inches:.2f}\", width={shape.width.inches:.2f}\", height={shape.height.inches:.2f}\"")
            if shape.has_text_frame:
                print(f"    Content: {text.strip().replace('\n', ' | ')}")
        print()

if __name__ == "__main__":
    inspect_shapes()
