#!/usr/bin/env python3
"""
Refined SIH Official Presentation Generator for SafeRoute AI
-------------------------------------------------------------
Populates `IDEA_Presentation_Format.pptx` by directly formatting official slide titles,
clearing placeholder prompt text, and rendering dark theme cards and diagrams.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Define Dark Theme Palette
COLOR_BG = RGBColor(11, 15, 25)          # Obsidian Dark #0B0F19
COLOR_CARD = RGBColor(30, 41, 59)        # Slate Dark Card #1E293B
COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate 700 Border #334155
COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# White Main Text #F8FAFC
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# Muted Gray Text #94A3B8
COLOR_ACCENT_CYAN = RGBColor(56, 189, 248)# Sky Cyan Accent #38BDF8
COLOR_ACCENT_EMERALD = RGBColor(52, 211, 153)# Safe Emerald Green #34D399
COLOR_ACCENT_RED = RGBColor(248, 113, 113)  # Hazard Red #F87171
COLOR_ACCENT_YELLOW = RGBColor(251, 191, 36) # Reroute Yellow #FBBF24
COLOR_ACCENT_PURPLE = RGBColor(167, 139, 250)# AI/Engine Purple #A78BFA
COLOR_CONTAINER_DARK = RGBColor(15, 23, 42) # Deep Slate Container #0F172A

def set_slide_background(slide, color=COLOR_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def set_speaker_notes(slide, script, points, transition, qa):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    notes_text = f"""=== SCRIPT / WHAT TO SAY ===
{script}

=== KEY TECHNICAL POINTS ===
{points}

=== SLIDE TRANSITION ===
{transition}

=== EXPECTED JURY Q&A ===
{qa}"""
    text_frame.text = notes_text

def clean_existing_content_shapes(slide):
    """
    Clears body placeholder/textbox shapes except footer and slide number.
    """
    to_remove = []
    for shape in slide.shapes:
        name = shape.name.lower()
        if "footer" in name or "slide number" in name:
            # Keep footer and slide number placeholders!
            continue
        to_remove.append(shape)
        
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

def add_header(slide, title_text):
    # Header Category / Title
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.6))
    tf = t_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.font.name = "Arial"

def build_sih_presentation():
    template_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/IDEA_Presentation_Format.pptx"
    prs = Presentation(template_path)
    
    # Remove slide 1 (Instructions Slide)
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]
    
    print(f"Instructions slide removed. Total slides remaining: {len(prs.slides)}")
    
    # SLIDE 1: TITLE PAGE
    print("Populating Slide 1: TITLE PAGE...")
    slide1 = prs.slides[0]
    set_slide_background(slide1)
    clean_existing_content_shapes(slide1)
    add_header(slide1, "TITLE PAGE")
    
    create_card(slide1, Inches(0.8), Inches(1.1), Inches(11.733), Inches(5.6), COLOR_CONTAINER_DARK, COLOR_CARD_BORDER)
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(11.733), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT_CYAN
    bar.line.fill.background()
    
    tb1 = slide1.shapes.add_textbox(Inches(1.1), Inches(1.3), Inches(11.1), Inches(5.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "SafeRoute AI — Dynamic Fire Evacuation Router"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "Decentralized Edge Mesh with Real-Time Physical Hazard Mapping & 3D Digital Twin"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_ACCENT_CYAN
    
    meta_items = [
        ("Problem Statement ID", "SIH-2026-SR01 (Dynamic Fire Evacuation Router)"),
        ("Problem Statement Title", "Dynamic Fire Evacuation Router with Real-Time Hazard Mapping"),
        ("Theme", "Smart Life-Safety Infrastructure & Edge AI Systems"),
        ("PS Category", "Hardware & Embedded Software Integration"),
        ("Team Name / Leads", "SafeRoute AI Core Engineering Team"),
        ("Key Innovations", "Sub-300ms Edge Dijkstra • ESP-NOW Mesh • 3-Tier Fail-Safe • 3D EOC Twin")
    ]
    for k, v in meta_items:
        pk = tf1.add_paragraph()
        pk.text = f"{k}: "
        pk.font.size = Pt(11)
        pk.font.bold = True
        pk.font.color.rgb = COLOR_ACCENT_YELLOW
        pv = pk.add_run()
        pv.text = v
        pv.font.size = Pt(11)
        pv.font.bold = False
        pv.font.color.rgb = COLOR_TEXT_MAIN
        
    set_speaker_notes(
        slide1,
        script="Good day judges. We present SafeRoute AI for Problem Statement ID SIH-2026-SR01: Dynamic Fire Evacuation Router with Real-Time Hazard Mapping. SafeRoute AI replaces passive exit signs with decentralized ESP32 edge micro-routers that compute safest exit paths in under 300 milliseconds without central server dependency.",
        points="• Official SIH Title Page format satisfying all metadata fields.\n• Highlight Problem Statement ID, Category, and core engineering pillars.",
        transition="Let's detail our proposed solution, prototype, and key innovations.",
        qa="Q: Is this a software or hardware submission?\nA: It is a combined Hardware & Software implementation: ESP32 micro-routers with WS2812B LEDs, paired with a Python simulator and 3D WebGL EOC dashboard."
    )

    # SLIDE 2: PROPOSED SOLUTION
    print("Populating Slide 2: PROPOSED SOLUTION...")
    slide2 = prs.slides[1]
    set_slide_background(slide2)
    clean_existing_content_shapes(slide2)
    add_header(slide2, "PROPOSED SOLUTION")
    
    cols = [
        ("PROPOSED SOLUTION / PROTOTYPE", [
            ("Decentralized Edge Graph", "Each node acts as an OSPF-style physical router running Dijkstra on-device."),
            ("Multi-Vector Sensor Fusion", "Fuses Temperature (°C), Smoke (PPM), Flame (IR), and Occupancy into dynamic edge weights."),
            ("Dynamic Visual Guidance", "Drives WS2812B chasing LED strips to direct occupants away from hazards.")
        ], COLOR_ACCENT_CYAN),
        
        ("HOW IT ADDRESSES THE PROBLEM", [
            ("Eliminates Static Trap", "Recalculates paths instantly when fire blocks an exit, reversing chasing LED direction."),
            ("Sub-300ms Reaction Time", "Updates visual indicators before occupant panic commits people to danger zones."),
            ("Zero Single-Point-of-Failure", "Runs 100% autonomously without reliance on central servers or power mains.")
        ], COLOR_ACCENT_EMERALD),
        
        ("INNOVATION & UNIQUENESS", [
            ("Continuous Exponential Math", "Uses exponential cost functions exp(αT+βS), avoiding crude binary thresholds."),
            ("Lock-Free FreeRTOS Core", "Core 0/1 double-buffering avoids mutex priority inversion on hot routing path."),
            ("3-Tier Fail-Safe Hierarchy", "Seamlessly transitions from local sensors to neighbor consensus on hardware fault.")
        ], COLOR_ACCENT_PURPLE)
    ]
    
    for i, (cat, items, color) in enumerate(cols):
        x = 0.8 + i * 3.98
        create_card(slide2, Inches(x), Inches(1.2), Inches(3.75), Inches(5.5), COLOR_CARD, COLOR_CARD_BORDER)
        hbar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.2), Inches(3.75), Inches(0.45))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = color
        hbar.line.fill.background()
        
        tb_cat = slide2.shapes.add_textbox(Inches(x + 0.1), Inches(1.25), Inches(3.55), Inches(0.35))
        tf_cat = tb_cat.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = cat
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CONTAINER_DARK
        
        tb_body = slide2.shapes.add_textbox(Inches(x + 0.15), Inches(1.75), Inches(3.45), Inches(4.85))
        tf_body = tb_body.text_frame
        tf_body.word_wrap = True
        for k, v in items:
            pk = tf_body.add_paragraph()
            pk.text = k
            pk.font.size = Pt(10.5)
            pk.font.bold = True
            pk.font.color.rgb = COLOR_TEXT_MAIN
            pv = tf_body.add_paragraph()
            pv.text = v
            pv.font.size = Pt(9.5)
            pv.font.color.rgb = COLOR_TEXT_MUTED
            
    set_speaker_notes(
        slide2,
        script="Slide 2 presents our proposed solution. SafeRoute AI replaces passive exit signs with decentralized link-state micro-routers. Each node senses multi-vector fire data and recomputes paths in under 300 milliseconds. Our innovation lies in our continuous exponential math formula, lock-free FreeRTOS double-buffering, and 3-tier fail-safe hierarchy.",
        points="• Strictly addresses SIH Slide 2 pointers: Proposed Solution, Problem Solved, Innovation & Uniqueness.\n• Highlight sub-300ms latency, continuous exponential cost formula, and lock-free execution.",
        transition="Let's now examine our technical approach, technologies used, and system flow.",
        qa="Q: What makes your algorithm unique compared to standard A*?\nA: Standard A* uses static heuristic distances. SafeRoute AI computes dynamic physical hazard weights continuously, combining thermal, particulate, and congestion vectors."
    )

    # SLIDE 3: TECHNICAL APPROACH
    print("Populating Slide 3: TECHNICAL APPROACH...")
    slide3 = prs.slides[2]
    set_slide_background(slide3)
    clean_existing_content_shapes(slide3)
    add_header(slide3, "TECHNICAL APPROACH")
    
    # Left Box: Technologies Used
    create_card(slide3, Inches(0.8), Inches(1.2), Inches(5.75), Inches(5.5), COLOR_CARD, COLOR_CARD_BORDER)
    tb_tech = slide3.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(5.35), Inches(5.1))
    tf_tech = tb_tech.text_frame
    tf_tech.word_wrap = True
    
    pt1 = tf_tech.paragraphs[0]
    pt1.text = "TECHNOLOGIES & ARCHITECTURE"
    pt1.font.size = Pt(11)
    pt1.font.bold = True
    pt1.font.color.rgb = COLOR_ACCENT_CYAN
    
    tech_items = [
        ("Embedded Hardware", "ESP32 Dual-Core 240MHz, DHT22 (Temp), MQ-2 (Smoke), IR Flame, WS2812B LEDs, Buzzer."),
        ("Firmware Stack", "C++17 under PlatformIO / FreeRTOS. Core 1 handles ESP-NOW, Core 0 runs Dijkstra."),
        ("Networking Mesh", "ESP-NOW 2.4GHz connectionless flooding with 24-byte packet, CRC16, and monotonic seq-nums."),
        ("Simulation & Twin", "Python digital twin injector tool, FastAPI backend bridge, Three.js 3D WebGL EOC UI.")
    ]
    for title, desc in tech_items:
        pk = tf_tech.add_paragraph()
        pk.text = title
        pk.font.size = Pt(10.5)
        pk.font.bold = True
        pk.font.color.rgb = COLOR_TEXT_MAIN
        pv = tf_tech.add_paragraph()
        pv.text = desc
        pv.font.size = Pt(9)
        pv.font.color.rgb = COLOR_TEXT_MUTED
        
    # Right Box: Math Formula & LED State Workflow
    create_card(slide3, Inches(6.8), Inches(1.2), Inches(5.733), Inches(5.5), COLOR_CARD, COLOR_CARD_BORDER)
    tb_math = slide3.shapes.add_textbox(Inches(7.0), Inches(1.35), Inches(5.333), Inches(5.1))
    tf_math = tb_math.text_frame
    tf_math.word_wrap = True
    
    pm1 = tf_math.paragraphs[0]
    pm1.text = "MATHEMATICAL FORMULA & LED MATRIX"
    pm1.font.size = Pt(11)
    pm1.font.bold = True
    pm1.font.color.rgb = COLOR_ACCENT_YELLOW
    
    form_text = "edge_cost = (base_dist · exp(2.2·T_norm + 1.6·S_norm) + 0.5·O_norm·base_dist) · (Flame ? 10⁶ : 1)"
    pm2 = tf_math.add_paragraph()
    pm2.text = form_text
    pm2.font.size = Pt(9.5)
    pm2.font.name = "Courier New"
    pm2.font.bold = True
    pm2.font.color.rgb = COLOR_TEXT_MAIN
    
    led_matrix = [
        ("🟢 GREEN", "Safe shortest path; chasing lights point to exit."),
        ("🟡 YELLOW", "Rerouted alternate route (Smoke-dominant: S_norm > T_norm)."),
        ("🔴 PULSING RED", "Immediate danger (Active flame OR T_norm ≥ S_norm)."),
        ("⚪ WHITE STROBE", "Shelter-in-place state (All exit costs ≥ 100,000).")
    ]
    for state, desc in led_matrix:
        pk = tf_math.add_paragraph()
        pk.text = state
        pk.font.size = Pt(10)
        pk.font.bold = True
        pk.font.color.rgb = COLOR_TEXT_MAIN
        pv = tf_math.add_paragraph()
        pv.text = desc
        pv.font.size = Pt(9)
        pv.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide3,
        script="Slide 3 outlines our technical approach. Our hardware stack uses dual-core ESP32 microcontrollers reading multi-vector sensors. Firmware is written in C++17 on PlatformIO/FreeRTOS, isolating ESP-NOW mesh reception to Core 1 and Dijkstra routing to Core 0. The physical cost math combines temperature, smoke, occupancy, and flame into an exponential edge weight calibrated via NIST fire datasets. FastLED states dictate visual evacuation guidance.",
        points="• Strictly satisfies SIH Slide 3 pointers: Technologies Used & Implementation Methodology.\n• Detail FreeRTOS core split, ESP-NOW 24B wire packet, continuous cost math, and FastLED state matrix.",
        transition="Let's evaluate the feasibility, risks, and mitigation strategies of our project.",
        qa="Q: How do you prevent LED path flickering?\nA: We implement a 1500–2000ms hold-down hysteresis timer that prevents non-essential route switching unless flame is detected."
    )

    # SLIDE 4: FEASIBILITY AND VIABILITY
    print("Populating Slide 4: FEASIBILITY AND VIABILITY...")
    slide4 = prs.slides[3]
    set_slide_background(slide4)
    clean_existing_content_shapes(slide4)
    add_header(slide4, "FEASIBILITY AND VIABILITY")
    
    fv_cols = [
        ("FEASIBILITY ANALYSIS", [
            ("Hardware Cost ($16-$24)", "Commodity ESP32 nodes and standard sensors keep per-node BOM extremely low."),
            ("Retrofit Friendly", "Installs alongside existing code-mandated static signage without rewiring."),
            ("Zero Cloud Cost", "Operates 100% autonomously without requiring expensive cloud subscriptions.")
        ], COLOR_ACCENT_CYAN),
        
        ("POTENTIAL CHALLENGES & RISKS", [
            ("Sensor Failure / Stuck ADC", "Physical sensors can break, short-circuit, or freeze ADC values during fire conditions."),
            ("Packet Corruption / Drop", "RF interference or fire damage can drop mesh packets or corrupt data payloads."),
            ("All Exit Traversal Blocked", "Severe flashover can block all available building exit paths simultaneously.")
        ], COLOR_ACCENT_RED),
        
        ("MITIGATION STRATEGIES", [
            ("3-Tier Fail-Safe Hierarchy", "Local variance check detects stuck sensors; automatically defers to Tier 2 Neighbor Consensus."),
            ("CRC16 & Monotonic Seq-Num", "Validates payload integrity and rejects replayed or corrupted packets live."),
            ("Shelter-In-Place State", "Detects cost ≥ 100,000, triggering White Strobe LEDs to instruct room sealing.")
        ], COLOR_ACCENT_EMERALD)
    ]
    
    for i, (cat, items, color) in enumerate(fv_cols):
        x = 0.8 + i * 3.98
        create_card(slide4, Inches(x), Inches(1.2), Inches(3.75), Inches(5.5), COLOR_CARD, COLOR_CARD_BORDER)
        hbar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.2), Inches(3.75), Inches(0.45))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = color
        hbar.line.fill.background()
        
        tb_cat = slide4.shapes.add_textbox(Inches(x + 0.1), Inches(1.25), Inches(3.55), Inches(0.35))
        tf_cat = tb_cat.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = cat
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CONTAINER_DARK
        
        tb_body = slide4.shapes.add_textbox(Inches(x + 0.15), Inches(1.75), Inches(3.45), Inches(4.85))
        tf_body = tb_body.text_frame
        tf_body.word_wrap = True
        for k, v in items:
            pk = tf_body.add_paragraph()
            pk.text = k
            pk.font.size = Pt(10.5)
            pk.font.bold = True
            pk.font.color.rgb = COLOR_TEXT_MAIN
            pv = tf_body.add_paragraph()
            pv.text = v
            pv.font.size = Pt(9.5)
            pv.font.color.rgb = COLOR_TEXT_MUTED
            
    set_speaker_notes(
        slide4,
        script="Slide 4 covers Feasibility and Viability. SafeRoute AI is highly feasible with a hardware cost of $16 to $24 per node, retrofitting easily alongside static exit signs. We identified three primary risks: sensor hardware failure, packet corruption, and total exit blockage. We mitigated all three: first, a 3-tier fail-safe hierarchy; second, CRC16 checksums and sequence numbers; and third, an automatic Shelter-In-Place state when all exits are blocked.",
        points="• Strictly satisfies SIH Slide 4 pointers: Feasibility Analysis, Challenges & Risks, Mitigation Strategies.\n• Detail 3-tier fail-safe, CRC16 validation, and Shelter-In-Place trigger.",
        transition="Let's inspect our project artifacts, code structure, and live dashboard interfaces.",
        qa="Q: What happens if a node's local sensor dies completely?\nA: The node detects zero sample variance over 30s, flags itself unhealthy, and transitions to Tier 2 (Neighbor Consensus) to compute physical hazard costs."
    )

    # SLIDE 5: ARTIFACTS
    print("Populating Slide 5: ARTIFACTS...")
    slide5 = prs.slides[4]
    set_slide_background(slide5)
    clean_existing_content_shapes(slide5)
    add_header(slide5, "ARTIFACTS")
    
    artifacts = [
        ("FIRMWARE & ALGORITHM CODE", [
            ("Modular Repository", "firmware/src/ (main, routing, fusion, comms, leds, failsafe)."),
            ("On-Device Dijkstra", "Non-blocking link-state graph solver running on ESP32 Core 0."),
            ("Double-Buffer Swap", "Lock-free atomic pointer swap synchronization.")
        ], COLOR_ACCENT_CYAN),
        
        ("FIRE INJECTOR SIMULATOR", [
            ("Multi-Profile Injector", "simulator/injector.py streams smolder & flashover profiles."),
            ("Judge Attack Trigger", "Allows judges to trigger flashover on arbitrary zones live."),
            ("Corrupt Packet Generator", "Injects malformed CRC payloads to prove live rejection.")
        ], COLOR_ACCENT_EMERALD),
        
        ("3D EOC DIGITAL TWIN", [
            ("Three.js WebGL Interface", "dashboard/ renders 2D/3D floor grid & exit paths live."),
            ("Async IDW Heatmap", "Inverse Distance Weighting interpolates heat/smoke surfaces."),
            ("Node Health Panel", "Surfaces sensor fault alerts and consensus state live.")
        ], COLOR_ACCENT_PURPLE),
        
        ("PHYSICAL TESTBENCH", [
            ("Hardware Prototypes", "ESP32 nodes wired with WS2812B strips, MQ-2, DHT22."),
            ("Wokwi Simulation", "Full hardware binary simulation suite under tests/."),
            ("Timing Benchmark Logs", "Measured p95 reaction latency of 118ms across 4 hops.")
        ], COLOR_ACCENT_YELLOW)
    ]
    
    for i, (cat, items, color) in enumerate(artifacts):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 5.95
        y = 1.2 + row * 2.75
        
        create_card(slide5, Inches(x), Inches(y), Inches(5.75), Inches(2.55), COLOR_CARD, COLOR_CARD_BORDER)
        tb = slide5.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(5.35), Inches(2.25))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = cat
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        for k, v in items:
            pk = tf.add_paragraph()
            pk.text = f"• {k}: "
            pk.font.size = Pt(10)
            pk.font.bold = True
            pk.font.color.rgb = COLOR_TEXT_MAIN
            pv = pk.add_run()
            pv.text = v
            pv.font.size = Pt(9.5)
            pv.font.bold = False
            pv.font.color.rgb = COLOR_TEXT_MUTED
            
    set_speaker_notes(
        slide5,
        script="Slide 5 displays our working project artifacts: modular C++ firmware, Python fire simulator supporting judge flashover triggers and CRC corruption tests, 3D WebGL EOC digital twin, and physical testbench benchmarks proving a 118ms p95 latency across 4 hops.",
        points="• Strictly satisfies SIH Slide 5 pointers: Embedded Code, Solution Snaps, Dashboard Snaps, Simulator.\n• Highlight physical testbench, Python injector, and Three.js 3D Digital Twin.",
        transition="Let's conclude with our research references and dataset calibrations.",
        qa="Q: Is the code open source and available for review?\nA: Yes, the full repository includes firmware source code, Python digital twin injector, Node-RED flows, tests, and Docker deployment scripts."
    )

    # SLIDE 6: RESEARCH AND REFERENCES
    print("Populating Slide 6: RESEARCH AND REFERENCES...")
    slide6 = prs.slides[5]
    set_slide_background(slide6)
    clean_existing_content_shapes(slide6)
    add_header(slide6, "RESEARCH AND REFERENCES")
    
    # Left Card: Public Datasets & Curve Regression
    create_card(slide6, Inches(0.8), Inches(1.2), Inches(5.75), Inches(5.5), COLOR_CARD, COLOR_CARD_BORDER)
    tb_ref1 = slide6.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(5.35), Inches(5.1))
    tf_ref1 = tb_ref1.text_frame
    tf_ref1.word_wrap = True
    
    pr1 = tf_ref1.paragraphs[0]
    pr1.text = "DATASETS & CURVE REGRESSION"
    pr1.font.size = Pt(11)
    pr1.font.bold = True
    pr1.font.color.rgb = COLOR_ACCENT_CYAN
    
    ds_items = [
        ("NIST Fire Dynamics Data", "Extracted temperature rise timelines and smoke flashover curves from National Institute of Standards and Technology public research."),
        ("Kaggle Smoke Datasets", "Used public smoke detection time-series data to fit logistic growth curves for smoldering fire profiles."),
        ("Hyperparameter Tuning", "Offline regression scripts under simulator/fire_profiles/ fit α=2.2 (thermal) and β=1.6 (smoke) parameters.")
    ]
    for title, desc in ds_items:
        pk = tf_ref1.add_paragraph()
        pk.text = title
        pk.font.size = Pt(10.5)
        pk.font.bold = True
        pk.font.color.rgb = COLOR_TEXT_MAIN
        pv = tf_ref1.add_paragraph()
        pv.text = desc
        pv.font.size = Pt(9)
        pv.font.color.rgb = COLOR_TEXT_MUTED
        
    # Right Card: Standards & Technical References
    create_card(slide6, Inches(6.8), Inches(1.2), Inches(5.733), Inches(5.5), COLOR_CARD, COLOR_CARD_BORDER)
    tb_ref2 = slide6.shapes.add_textbox(Inches(7.0), Inches(1.35), Inches(5.333), Inches(5.1))
    tf_ref2 = tb_ref2.text_frame
    tf_ref2.word_wrap = True
    
    pr2 = tf_ref2.paragraphs[0]
    pr2.text = "STANDARDS & TECHNICAL REFERENCES"
    pr2.font.size = Pt(11)
    pr2.font.bold = True
    pr2.font.color.rgb = COLOR_ACCENT_EMERALD
    
    st_items = [
        ("NFPA 101 Life Safety Code", "Framed system to augment — not replace — code-mandated static illuminated exit signage in commercial facilities."),
        ("OSPF RFC 2328 Protocol", "Applied classical link-state network graph routing concepts to physical hazard topology modeling."),
        ("FreeRTOS Dual-Core Architecture", "Implemented lock-free double-buffered pointer swapping based on real-time embedded systems literature.")
    ]
    for title, desc in st_items:
        pk = tf_ref2.add_paragraph()
        pk.text = title
        pk.font.size = Pt(10.5)
        pk.font.bold = True
        pk.font.color.rgb = COLOR_TEXT_MAIN
        pv = tf_ref2.add_paragraph()
        pv.text = desc
        pv.font.size = Pt(9)
        pv.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide6,
        script="Slide 6 presents our research foundation and references. Our cost function constants α=2.2 and β=1.6 were fitted using offline non-linear regressions on NIST fire dynamics data and Kaggle smoke datasets. SafeRoute AI is designed to augment static signage under NFPA 101 standards, applying classical OSPF link-state graph routing principles to physical fire evacuation. Thank you, and we are open for jury questions.",
        points="• Strictly satisfies SIH Slide 6 pointers: Reference Links & Research Work Details.\n• Highlight NIST fire data, Kaggle smoke time-series, NFPA 101 regulatory framing, and OSPF graph theory.",
        transition="End of presentation. Open floor for jury Q&A.",
        qa="Q: How do you justify fitting hyperparameters offline instead of running ML on-device?\nA: Safety-critical emergency routing requires deterministic execution. Offline regression calibrates accurate continuous physics curves while ensuring on-device Dijkstra runs in bounded sub-300ms time."
    )

    output_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/SafeRouteAI_SIH_Official_Submission.pptx"
    docs_output_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/docs/presentations/SafeRouteAI_SIH_Official_Submission.pptx"
    prs.save(output_path)
    prs.save(docs_output_path)
    print(f"Refined Official SIH Submission Presentation saved to:\n  - {output_path}\n  - {docs_output_path}")

if __name__ == "__main__":
    build_sih_presentation()
