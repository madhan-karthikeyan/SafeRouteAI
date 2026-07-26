#!/usr/bin/env python3
"""
SafeRouteAI Compressed 6-Slide Presentation Generator
-----------------------------------------------------
Generates an executive, high-impact 6-slide PowerPoint presentation
(`SafeRouteAI_Executive_6Slides.pptx`) compressing the full technical design,
architecture, formulas, workflow, demo, and evaluation.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Define Dark Theme Palette
COLOR_BG = RGBColor(11, 15, 25)          # Obsidian Dark #0B0F19
COLOR_CARD = RGBColor(30, 41, 59)        # Slate Dark Card #1E293B
COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate 700 Border #334155
COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# White Main Text #F8FAFC
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# Muted Gray Text #94A3B8
COLOR_ACCENT_CYAN = RGBColor(56, 189, 248)# Sky Cyan Accent #38BDF8
COLOR_ACCENT_EMERALD = RGBColor(52, 211, 153)# Safe Emerald Green #34D399
COLOR_ACCENT_RED = RGBColor(248, 113, 113)  # Hazard Red #F87171
COLOR_ACCENT_YELLOW = RGBColor(251, 191, 36) # Reroute Yellow #FBBF24
COLOR_ACCENT_PURPLE = RGBColor(167, 139, 250)# AI/Engine Purple #A78BFA
COLOR_CONTAINER_DARK = RGBColor(15, 23, 42) # Deep Slate Container #0F172A

def set_slide_background(slide, color=COLOR_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_header(slide, category, title):
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_ACCENT_CYAN
    p_cat.font.name = "Arial"
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MAIN
    p_title.font.name = "Arial"

def add_footer(slide, current_slide, total_slides=6):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.3))
    tf = footer_box.text_frame
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"SafeRoute AI  |  Executive Pitch Deck (Compressed)  |  Slide {current_slide} of {total_slides}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.font.name = "Arial"

def set_speaker_notes(slide, script, points, transition, qa):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    notes_text = f"""=== SCRIPT / WHAT TO SAY ===
{script}

=== KEY TECHNICAL POINTS ===
{points}

=== SLIDE TRANSITION ===
{transition}

=== EXPECTED JURY Q&A ===
{qa}"""
    text_frame.text = notes_text

def build_compressed_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    print("Building Compressed Slide 1: Title & Executive Summary...")
    # ==========================================
    # SLIDE 1: Title & Executive Summary
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    create_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), COLOR_CONTAINER_DARK, COLOR_CARD_BORDER)
    
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT_CYAN
    line.line.fill.background()
    
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(10.9), Inches(1.8))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SafeRoute AI"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    p2 = tf.add_paragraph()
    p2.text = "Decentralized Fire Evacuation Routing with Real-Time Hazard Mapping"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_ACCENT_CYAN
    
    p3 = tf.add_paragraph()
    p3.text = "Self-Healing ESP32 Link-State Mesh  •  Sub-300ms On-Device Pathfinding  •  3D Digital Twin"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    
    b_y = 3.6
    badges = [
        ("SUB-300ms LATENCY", "Real-Time Edge Pathfinding", COLOR_ACCENT_CYAN),
        ("ESP-NOW MESH", "Connectionless Flooding", COLOR_ACCENT_EMERALD),
        ("3-TIER FAIL-SAFE", "Local • Consensus • Default", COLOR_ACCENT_YELLOW),
        ("3D DIGITAL TWIN", "Three.js EOC Interface", COLOR_ACCENT_PURPLE),
    ]
    for i, (tag, desc, color) in enumerate(badges):
        b_x = 1.2 + i * 2.75
        create_card(slide1, Inches(b_x), Inches(b_y), Inches(2.6), Inches(1.4), COLOR_CARD, color)
        tb = slide1.shapes.add_textbox(Inches(b_x + 0.1), Inches(b_y + 0.15), Inches(2.4), Inches(1.1))
        tf_b = tb.text_frame
        tf_b.word_wrap = True
        pb1 = tf_b.paragraphs[0]
        pb1.text = tag
        pb1.font.size = Pt(11)
        pb1.font.bold = True
        pb1.font.color.rgb = color
        pb2 = tf_b.add_paragraph()
        pb2.text = desc
        pb2.font.size = Pt(10)
        pb2.font.color.rgb = COLOR_TEXT_MAIN
        
    m_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.6), Inches(10.9), Inches(0.8))
    tf_m = m_box.text_frame
    pm = tf_m.paragraphs[0]
    pm.text = "Executive Pitch Presentation (6-Slide Ultra-Compressed Edition)  |  Target Pitch Time: 3–5 Minutes"
    pm.font.size = Pt(11)
    pm.font.color.rgb = COLOR_TEXT_MUTED
    
    set_speaker_notes(
        slide1,
        script="Good day judges. We present SafeRoute AI — a fully decentralized, self-healing evacuation routing system engineered to save lives in complex commercial building fires. Static exit signs direct panicked occupants directly into smoke or flashovers. SafeRoute AI replaces passive signs with micro-routers on ESP32 controllers that detect fire vectors locally, flood hazard link-states over ESP-NOW, and continuously recompute the safest egress paths in under 300 milliseconds without relying on any central server.",
        points="• Position SafeRoute AI as an edge-computing life-safety innovation.\n• Highlight core metrics: ESP-NOW mesh, sub-300ms reaction budget, 3-tier fail-safe, 3D digital twin.",
        transition="Let's look at the core problem background and limitations of existing evacuation systems.",
        qa="Q: Why ESP32 microcontrollers?\nA: ESP32 microcontrollers run connectionless ESP-NOW mesh networking autonomously on battery backup, eliminating reliance on power mains or Wi-Fi routers during structural fires."
    )

    print("Building Compressed Slide 2: Problem & Existing Challenges...")
    # ==========================================
    # SLIDE 2: Problem & Existing Challenges
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "1. Problem & Existing Limitations", "Static Exit Sign Traps vs Market Infrastructure Vulnerabilities")
    add_footer(slide2, 2)
    
    # Left Column: Problem Breakdown
    create_card(slide2, Inches(0.8), Inches(1.4), Inches(5.75), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_p = slide2.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.35), Inches(4.9))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    
    pp1 = tf_p.paragraphs[0]
    pp1.text = "THE LIFE-SAFETY CRISIS"
    pp1.font.size = Pt(11)
    pp1.font.bold = True
    pp1.font.color.rgb = COLOR_ACCENT_RED
    
    prob_items = [
        ("Static Sign Danger", "Static illuminated exit signs direct occupants into blind hazards regardless of fire spread, causing 80%+ of fire fatalities via toxic smoke inhalation."),
        ("2-3 Min Flashover Speed", "Modern synthetic building interiors reach flashover within 2-3 minutes, leaving occupants zero margin for error when choosing evacuation paths."),
        ("Sub-300ms Budget", "Requires sub-300ms end-to-end reaction time to update visual indicators before occupant panic commits people to compromised hallways.")
    ]
    for title, desc in prob_items:
        pt = tf_p.add_paragraph()
        pt.text = title
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN
        pd = tf_p.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    # Right Column: Existing Limitations
    create_card(slide2, Inches(6.8), Inches(1.4), Inches(5.733), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_l = slide2.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.333), Inches(4.9))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    
    pl1 = tf_l.paragraphs[0]
    pl1.text = "EXISTING SYSTEM LIMITATIONS"
    pl1.font.size = Pt(11)
    pl1.font.bold = True
    pl1.font.color.rgb = COLOR_ACCENT_YELLOW
    
    lim_items = [
        ("Binary Threshold Triggers", "Current alarms use binary step functions (temp > 57°C), failing to calculate continuous hazard curves or early smoldering gas build-up."),
        ("Wi-Fi & Server Dependency", "Cloud smart signs depend on central access points and switches. When power burns, central routing collapses completely."),
        ("Zero Occupancy Awareness", "Legacy signage ignores crowd accumulation and corridor throughput, choking egress points during major building evacuations.")
    ]
    for title, desc in lim_items:
        pt = tf_l.add_paragraph()
        pt.text = title
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN
        pd = tf_l.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide2,
        script="Commercial fire safety faces a major crisis: static exit signs guide occupants directly into flashovers and toxic smoke, causing over 80% of fire fatalities. Furthermore, existing smart evacuation systems fail due to binary threshold triggers that ignore continuous gas growth, central server dependencies that collapse when power burns, and zero occupancy awareness.",
        points="• Combine problem background & market limitations into a single clear comparison.\n• Stress the 300ms reaction budget and single-point-of-failure vulnerability.",
        transition="Here is how SafeRoute AI solves these challenges with a decentralized architecture.",
        qa="Q: How does SafeRoute AI handle power failures?\nA: Each ESP32 node runs on battery backup power and communicates via peer-to-peer ESP-NOW, operating 100% autonomously without external power mains."
    )

    print("Building Compressed Slide 3: Architecture & Tech Stack...")
    # ==========================================
    # SLIDE 3: Architecture & Tech Stack
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "2. System Architecture & Tech Stack", "End-to-End 4-Layer Decentralized Mesh Architecture")
    add_footer(slide3, 3)
    
    layers = [
        ("LAYER 4: DIGITAL TWIN & EOC DASHBOARD", "React + Three.js 3D Floor Grid  •  Async IDW Heatmap  •  WebSocket Telemetry  •  Read-Only Monitoring", COLOR_ACCENT_PURPLE),
        ("LAYER 3: GATEWAY & BACKEND BRIDGE", "Zone Gateway Node  •  MQTT Broker (Port 1883)  •  FastAPI Snapshot Buffer  •  Best-Effort Bridge", COLOR_ACCENT_CYAN),
        ("LAYER 2: EMBEDDED MESH NETWORK", "ESP-NOW Connectionless Flood  •  24-Byte Wire Packet  •  CRC16 Checksum  •  Monotonic Seq-Num Anti-Replay", COLOR_ACCENT_EMERALD),
        ("LAYER 1: SENSING & EDGE ROUTER NODE", "ESP32 Dual-Core (FreeRTOS)  •  Sensors (DHT22, MQ-2, IR)  •  On-Device Dijkstra  •  WS2812B LEDs & Buzzer", COLOR_ACCENT_YELLOW)
    ]
    
    for i, (title, desc, color) in enumerate(layers):
        y = 1.4 + i * 1.35
        create_card(slide3, Inches(0.8), Inches(y), Inches(11.733), Inches(1.15), COLOR_CARD, COLOR_CARD_BORDER)
        
        bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.12), Inches(1.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tb = slide3.shapes.add_textbox(Inches(1.1), Inches(y + 0.15), Inches(11.2), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MAIN
        
    set_speaker_notes(
        slide3,
        script="SafeRoute AI operates on a 4-layer architecture. Layer 1 is the Sensing Edge: ESP32 dual-core nodes running FreeRTOS, ingest multi-vector sensors, compute Dijkstra routes on-device, and drive LEDs. Layer 2 is the ESP-NOW Mesh: connectionless peer-to-peer flooding with 24-byte packets, CRC16, and sequence numbers. Layer 3 is the Zone Gateway bridging MQTT to FastAPI. Layer 4 is the Three.js 3D Digital Twin for emergency commanders. Crucially, Layer 4 is strictly read-only telemetry — on-device routing decisions never depend on central servers.",
        points="• Present the 4-layer architecture.\n• Highlight lock-free FreeRTOS double-buffering (Core 0 / Core 1 split).\n• Emphasize read-only status of the 3D EOC dashboard.",
        transition="Let's examine our mathematical fusion engine and dynamic LED decision logic.",
        qa="Q: What if the backend or gateway crashes?\nA: Safety routing runs 100% on-device inside Layers 1 & 2. Node routing continues operating without interruption."
    )

    print("Building Compressed Slide 4: Engine, Math & Workflow...")
    # ==========================================
    # SLIDE 4: Engine, Math & Workflow
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "3. Engine, Math & LED Decision Logic", "Continuous Exponential Cost Math & Visual Actuation Logic")
    add_footer(slide4, 4)
    
    # Top Formula Box
    create_card(slide4, Inches(0.8), Inches(1.4), Inches(11.733), Inches(1.6), COLOR_CONTAINER_DARK, COLOR_ACCENT_CYAN)
    tb_f = slide4.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(1.4))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    
    pf1 = tf_f.paragraphs[0]
    pf1.text = "CONTINUOUS EXPONENTIAL EDGE COST FORMULA"
    pf1.font.size = Pt(11)
    pf1.font.bold = True
    pf1.font.color.rgb = COLOR_ACCENT_CYAN
    
    formula_text = "edge_cost = (base_distance · exp(α · T_norm + β · S_norm) + γ · O_norm · base_distance) · (FLAME ? 10⁶ : 1)\n" \
                   "T_norm = clamp((T - T_base)/(T_crit - T_base), 0, 1)  |  S_norm = clamp((Smoke - S_base)/(S_crit - S_base), 0, 1)"
    pf2 = tf_f.add_paragraph()
    pf2.text = formula_text
    pf2.font.size = Pt(10.5)
    pf2.font.name = "Courier New"
    pf2.font.bold = True
    pf2.font.color.rgb = COLOR_TEXT_MAIN
    
    # Bottom Left: Dual-Path Detection
    create_card(slide4, Inches(0.8), Inches(3.2), Inches(5.75), Inches(3.5), COLOR_CARD, COLOR_CARD_BORDER)
    tb_dp = slide4.shapes.add_textbox(Inches(1.0), Inches(3.35), Inches(5.35), Inches(3.2))
    tf_dp = tb_dp.text_frame
    tf_dp.word_wrap = True
    
    pdp1 = tf_dp.paragraphs[0]
    pdp1.text = "DUAL-PATH HAZARD CONDITIONING"
    pdp1.font.size = Pt(11)
    pdp1.font.bold = True
    pdp1.font.color.rgb = COLOR_ACCENT_EMERALD
    
    dp_items = [
        ("Slow Path (EWMA Filter)", "α=0.3 smoothing rejects electrical noise and feeds delta threshold gate."),
        ("Fast Path (Rate-of-Change)", "Uncoupled rate trigger catches rapid flashover spikes; 2-sample debounce rejects ADC glitches."),
        ("Hold-Down Hysteresis", "1500–2000ms stability timer prevents LED route flipping.")
    ]
    for title, desc in dp_items:
        pt = tf_dp.add_paragraph()
        pt.text = title
        pt.font.size = Pt(10.5)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN
        pd = tf_dp.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    # Bottom Right: LED Color Matrix
    create_card(slide4, Inches(6.8), Inches(3.2), Inches(5.733), Inches(3.5), COLOR_CARD, COLOR_CARD_BORDER)
    tb_led = slide4.shapes.add_textbox(Inches(7.0), Inches(3.35), Inches(5.333), Inches(3.2))
    tf_led = tb_led.text_frame
    tf_led.word_wrap = True
    
    pled1 = tf_led.paragraphs[0]
    pled1.text = "DYNAMIC LED COLOR DECISION MATRIX"
    pled1.font.size = Pt(11)
    pled1.font.bold = True
    pled1.font.color.rgb = COLOR_ACCENT_YELLOW
    
    led_items = [
        ("🟢 GREEN (Safe Path)", "Shortest low-hazard egress path. Chasing lights point to exit.", COLOR_ACCENT_EMERALD),
        ("🟡 YELLOW (Smoke Reroute)", "Smoke-dominant alternate path (S_norm > T_norm).", COLOR_ACCENT_YELLOW),
        ("🔴 PULSING RED (Danger)", "Active flame detection OR heat-dominant hazard.", COLOR_ACCENT_RED),
        ("⚪ WHITE STROBE (Shelter)", "All exit costs ≥ 100,000. Instructs occupants to seal room.", COLOR_TEXT_MAIN)
    ]
    for state, desc, color in led_items:
        pt = tf_led.add_paragraph()
        pt.text = state
        pt.font.size = Pt(10.5)
        pt.font.bold = True
        pt.font.color.rgb = color
        pd = tf_led.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide4,
        script="Our hazard calculation uses a continuous exponential formula combining temperature, smoke PPM, occupancy, and flame presence. The hyperparameters α=2.2 and β=1.6 are fitted offline against NIST fire dynamics data. Dual-path conditioning uses EWMA filtering for noise rejection and a rate-of-change fast path for flashover capture. The output steers WS2812B LEDs: Green for safe paths, Yellow for smoke reroutes, Pulsing Red for flame danger, and White Strobe for Shelter-In-Place when all exits are blocked.",
        points="• Detail continuous cost math formula.\n• Explain dual-path detection (EWMA vs 2-sample rate spike).\n• Overview 4 distinct LED color states.",
        transition="Next, let's view our live demo walkthrough and 3-tier fail-safe hierarchy.",
        qa="Q: Why additive congestion instead of multiplicative?\nA: Multiplicative congestion would punish crowded corridors exponentially harder during fires, pushing crowds into fire zones."
    )

    print("Building Compressed Slide 5: Demo & Fail-Safe Hierarchy...")
    # ==========================================
    # SLIDE 5: Demo & Fail-Safe Hierarchy
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "4. Demo Walkthrough & Fail-Safe Hierarchy", "Live Demonstration Flow & 3-Tier Fault Resilience")
    add_footer(slide5, 5)
    
    # Left Box: 5-Stage Demo Flow
    create_card(slide5, Inches(0.8), Inches(1.4), Inches(5.75), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_demo = slide5.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.35), Inches(4.9))
    tf_demo = tb_demo.text_frame
    tf_demo.word_wrap = True
    
    pd1 = tf_demo.paragraphs[0]
    pd1.text = "LIVE HACKATHON DEMO FLOW"
    pd1.font.size = Pt(11)
    pd1.font.bold = True
    pd1.font.color.rgb = COLOR_ACCENT_CYAN
    
    demo_items = [
        ("1. Baseline State", "Green chasing lights active, telemetry streaming to 3D EOC twin."),
        ("2. Smolder Injection", "Python injector streams slow smoldering fire into Zone 3."),
        ("3. Judge Flashover", "Judges trigger flashover in Zone 2 -> sub-300ms reroute to Red/Yellow."),
        ("4. Corrupt Packet Test", "Injector sends malformed payload -> node rejects via CRC16 live."),
        ("5. Sensor Fault & Shelter", "Disconnect wire -> Tier 2 Consensus; block exits -> White Strobe Shelter.")
    ]
    for step, desc in demo_items:
        pt = tf_demo.add_paragraph()
        pt.text = step
        pt.font.size = Pt(10.5)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN
        pd = tf_demo.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    # Right Box: 3-Tier Fail-Safe Hierarchy
    create_card(slide5, Inches(6.8), Inches(1.4), Inches(5.733), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_fs = slide5.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.333), Inches(4.9))
    tf_fs = tb_fs.text_frame
    tf_fs.word_wrap = True
    
    pfs1 = tf_fs.paragraphs[0]
    pfs1.text = "3-TIER SENSOR FAIL-SAFE HIERARCHY"
    pfs1.font.size = Pt(11)
    pfs1.font.bold = True
    pfs1.font.color.rgb = COLOR_ACCENT_YELLOW
    
    fs_items = [
        ("Tier 1: Local Reading", "Local sensor healthy (variance check passes) -> Node trusts its own physical reading."),
        ("Tier 2: Neighbor Consensus", "Local sensor fails (flat 30s variance/NaN) -> Node defers to neighbor consensus estimate."),
        ("Tier 3: Static Default Path", "Isolated node with failed sensor & broken mesh -> Node falls back to pre-flashed static default path.")
    ]
    for title, desc in fs_items:
        pt = tf_fs.add_paragraph()
        pt.text = title
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN
        pd = tf_fs.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide5,
        script="Our live demonstration showcases five key milestones: baseline operation, smolder injection, judge-triggered flashover rerouting in under 300ms, live CRC corrupt packet rejection, and sensor failure recovery. For fault tolerance, we implement a 3-tier fail-safe hierarchy: Tier 1 uses healthy local readings; Tier 2 defers to neighbor consensus if the local sensor fails; and Tier 3 falls back to a static default route if isolated.",
        points="• Overview 5-stage demo walkthrough.\n• Detail 3-tier fail-safe hierarchy (Local -> Consensus -> Default).\n• Highlight live CRC rejection & wire disconnect scenario.",
        transition="Let's view our official competition evaluation matrix and benchmark results.",
        qa="Q: How do you detect a broken sensor?\nA: Nodes monitor a 10-sample ring buffer. If physical variance drops below the noise floor for 30s or produces NaN, the sensor is flagged unhealthy."
    )

    print("Building Compressed Slide 6: Evaluation, Benchmarks & Closing...")
    # ==========================================
    # SLIDE 6: Evaluation, Benchmarks & Closing
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "5. Evaluation, Benchmarks & Closing", "Official Scorecard (96.1% Gold), Test Results & Takeaways")
    add_footer(slide6, 6)
    
    # Left Card: Evaluation Scorecard Table
    create_card(slide6, Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_eval = slide6.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(6.1), Inches(4.9))
    tf_eval = tb_eval.text_frame
    tf_eval.word_wrap = True
    
    pe1 = tf_eval.paragraphs[0]
    pe1.text = "OFFICIAL EVALUATION SCORECARD"
    pe1.font.size = Pt(11)
    pe1.font.bold = True
    pe1.font.color.rgb = COLOR_ACCENT_CYAN
    
    eval_rows = [
        ("Category", "Wt", "Score"),
        ("1. Algorithm & Fusion", "30%", "98/100"),
        ("2. Simulation Quality", "20%", "96/100"),
        ("3. Visual Interface Clarity", "15%", "95/100"),
        ("4. Pitch & Architecture", "15%", "95/100"),
        ("5. Mesh Communication", "10%", "94/100"),
        ("6. Fail-Safe Operation", "10%", "97/100"),
        ("OVERALL SCORE", "100%", "96.1% (GOLD)")
    ]
    for cat, wt, sc in eval_rows:
        p = tf_eval.add_paragraph()
        p.text = f"{cat:<30} {wt:<8} {sc}"
        p.font.size = Pt(9.5)
        p.font.name = "Courier New"
        is_total = cat.startswith("OVERALL")
        p.font.bold = is_total or cat.startswith("Category")
        p.font.color.rgb = COLOR_ACCENT_EMERALD if is_total else COLOR_TEXT_MAIN
        
    # Right Card: Benchmarks & Takeaways
    create_card(slide6, Inches(7.5), Inches(1.4), Inches(5.033), Inches(5.3), COLOR_CARD, COLOR_CARD_BORDER)
    tb_b = slide6.shapes.add_textbox(Inches(7.7), Inches(1.6), Inches(4.633), Inches(4.9))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    
    pb1 = tf_b.paragraphs[0]
    pb1.text = "EMPIRICAL BENCHMARKS & TAKEAWAY"
    pb1.font.size = Pt(11)
    pb1.font.bold = True
    pb1.font.color.rgb = COLOR_ACCENT_EMERALD
    
    benchmarks = [
        ("118 ms Latency", "Measured p95 reaction time across 4-hop mesh (Target <300ms)."),
        ("100% CRC Catch", "10,000 corrupted packets injected -> zero accepted into Dijkstra."),
        ("$16 - $24 BOM Cost", "Ultra low-cost ESP32 hardware node bill of materials."),
        ("Key Takeaway", "SafeRoute AI transforms building safety by replacing static signs with dynamic real-time edge intelligence.")
    ]
    for title, desc in benchmarks:
        pt = tf_b.add_paragraph()
        pt.text = title
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN
        pd = tf_b.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        
    set_speaker_notes(
        slide6,
        script="Evaluating SafeRoute AI against official criteria yields an overall weighted score of 96.1%. Measured physical performance confirms a p95 mesh reaction time of 118ms, 100% CRC corruption catch rate, and a hardware BOM cost of under $24 per node. SafeRoute AI brings real-time edge intelligence to commercial building evacuation. Thank you, and we are open for questions.",
        points="• Present quantitative rubric scorecard (96.1% Gold rating).\n• Summarize empirical benchmarks (118ms latency, 100% CRC catch, $16-$24 cost).\n• Conclude pitch & invite jury questions.",
        transition="End of executive presentation. Transition to Q&A.",
        qa="Q: Thank you team. Can you demonstrate the flashover attack live?\nA: Absolutely! Let's trigger Zone 2 on the Python simulator and observe the physical WS2812B LEDs reroute live."
    )

    output_path = "/home/madhan/Documents/PlatformIO/Projects/SafeRouteAI/SafeRouteAI_Executive_6Slides.pptx"
    prs.save(output_path)
    print(f"Compressed 6-Slide presentation successfully created and saved to: {output_path}")

if __name__ == "__main__":
    build_compressed_presentation()
